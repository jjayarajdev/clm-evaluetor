"""Generate Evaluetor Capabilities Briefing — 6-slide PPTX matching the HTML slider.

Output: sales/Evaluetor-Capabilities-Briefing.pptx
Run:    cd sales && uv run python generate_prospect_slider.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Palette (matches HTML slider) ────────────────────────────────
PAPER      = RGBColor(0xF7, 0xF6, 0xF3)
PAPER_2    = RGBColor(0xEF, 0xED, 0xE8)
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x0E, 0x0E, 0x10)
INK_2      = RGBColor(0x2D, 0x2F, 0x36)
INK_3      = RGBColor(0x5B, 0x5F, 0x66)
INK_4      = RGBColor(0xA0, 0xA4, 0xAB)
RULE       = RGBColor(0x1A, 0x1A, 0x1C)
RULE_2     = RGBColor(0xDD, 0xDB, 0xD5)
ORANGE     = RGBColor(0xE9, 0x4E, 0x1B)
ORANGE_2   = RGBColor(0xC7, 0x3E, 0x0F)

# Fonts available on PowerPoint Windows/Mac default install
FONT_SANS  = "Calibri"           # Geist analogue
FONT_MONO  = "Consolas"          # Geist Mono analogue
FONT_ITAL  = "Cambria"           # Instrument Serif italic analogue

# Slide dimensions: 16:9
SW = Inches(13.333); SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]  # blank

# ─── Helpers ──────────────────────────────────────────────────────
def set_bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_width is not None: sh.line.width = line_width
    sh.shadow.inherit = False
    return sh

def hline(slide, x, y, w, color, thickness=Pt(0.75)):
    """A 1px-equivalent horizontal rule."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(1))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    return sh

def vline(slide, x, y, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(1), h)
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    return sh

def tx(slide, left, top, width, height, text,
       font=FONT_SANS, size=14, color=INK, bold=False, italic=False,
       align=PP_ALIGN.LEFT, spacing=0, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing

    if isinstance(text, list):
        # list of (text, opts) tuples for mixed-style runs
        for i, item in enumerate(text):
            if isinstance(item, tuple):
                t, opts = item
            else:
                t, opts = item, {}
            r = p.add_run() if i > 0 else p.runs[0] if p.runs else p.add_run()
            r.text = t
            f = r.font
            f.name = opts.get('font', font)
            f.size = Pt(opts.get('size', size))
            f.bold = opts.get('bold', bold)
            f.italic = opts.get('italic', italic)
            f.color.rgb = opts.get('color', color)
            if spacing:
                from pptx.oxml.ns import qn
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(spacing * 100)))
    else:
        p.text = text
        for r in p.runs:
            f = r.font
            f.name = font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
            if spacing:
                from pptx.oxml.ns import qn
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(spacing * 100)))
    return box

def slide_chrome(slide, head_label, foot_label, bg=PAPER):
    """Common header (brand + label) and footer (signature + slide marker)."""
    set_bg(slide, bg)

    # Brand mark + name (top-left)
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.28), Inches(0.28), ORANGE_2)
    tx(slide, Inches(0.55), Inches(0.43), Inches(0.28), Inches(0.28),
       "E", font=FONT_MONO, size=11, color=SURFACE, bold=True, align=PP_ALIGN.CENTER)
    tx(slide, Inches(0.92), Inches(0.42), Inches(2.5), Inches(0.3),
       "Evaluetor", font=FONT_SANS, size=14, color=INK, bold=True)

    # Head label (top-right)
    tx(slide, Inches(8.0), Inches(0.43), Inches(4.8), Inches(0.3),
       head_label, font=FONT_MONO, size=9, color=INK_3, spacing=1.6,
       align=PP_ALIGN.RIGHT)

    # Top hairline
    hline(slide, Inches(0.55), Inches(0.85), Inches(12.23), RULE)

    # Footer signature (bottom)
    hline(slide, Inches(0.55), Inches(6.95), Inches(12.23), RULE_2)
    tx(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.25),
       "Senority B.V. · Rotterdam, NL",
       font=FONT_MONO, size=9, color=INK_3, spacing=1.4)
    tx(slide, Inches(7.0), Inches(7.05), Inches(5.78), Inches(0.25),
       foot_label,
       font=FONT_MONO, size=9, color=INK_3, spacing=1.4, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 01 · COVER
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "COVER", "Slide 01 — Cover", bg=PAPER)

# Big orange square mark
rect(s, Inches(0.7), Inches(1.6), Inches(0.85), Inches(0.85), ORANGE_2)
tx(s, Inches(0.7), Inches(1.7), Inches(0.85), Inches(0.7),
   "E", font=FONT_MONO, size=32, color=SURFACE, bold=True, align=PP_ALIGN.CENTER)

# Wordmark
tx(s, Inches(0.7), Inches(2.7), Inches(10), Inches(1.5),
   "Evaluetor", font=FONT_SANS, size=80, color=INK, bold=False)

# Tagline (mixed: sans then italic serif accent)
tx(s, Inches(0.7), Inches(4.05), Inches(11), Inches(0.6),
   [("AI-native contract lifecycle management ", {"size": 26, "color": INK_2, "font": FONT_SANS}),
    ("for European enterprises.", {"size": 26, "color": ORANGE_2, "font": FONT_ITAL, "italic": True})])

# Audience line
tx(s, Inches(0.7), Inches(4.85), Inches(11), Inches(0.5),
   "Built for Legal, Procurement, Finance, and Operations leaders managing the post-signature value zone.",
   font=FONT_ITAL, size=16, color=INK_3, italic=True)

# Briefing metadata — 3 columns
hline(s, Inches(0.7), Inches(5.85), Inches(10), RULE_2)
tx(s, Inches(0.7), Inches(6.0), Inches(3), Inches(0.25),
   "BRIEFING", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(0.7), Inches(6.3), Inches(3), Inches(0.3),
   "Platform & capabilities", font=FONT_SANS, size=13, color=INK, bold=True)

tx(s, Inches(4.2), Inches(6.0), Inches(3), Inches(0.25),
   "PREPARED", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(4.2), Inches(6.3), Inches(3), Inches(0.3),
   "Q2 2026", font=FONT_SANS, size=13, color=INK, bold=True)

tx(s, Inches(7.7), Inches(6.0), Inches(4), Inches(0.25),
   "FROM", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(7.7), Inches(6.3), Inches(4), Inches(0.3),
   "Senority B.V. · Rotterdam, NL", font=FONT_SANS, size=13, color=INK, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 02 · PLATFORM AT A GLANCE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "02 / THE PLATFORM", "Slide 02 — The platform", bg=PAPER)

# Headline
tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Nine agents. One model. ", {"size": 36, "color": INK, "font": FONT_SANS, "bold": False}),
    ("Every contract.", {"size": 36, "color": INK_2, "font": FONT_ITAL, "italic": True})])

# Subtitle
tx(s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.4),
   "Six built-in capabilities. None of them an add-on. None sold separately.",
   font=FONT_SANS, size=14, color=INK_3)

# Stats bar — 4 columns, top + bottom rule
hline(s, Inches(0.55), Inches(2.5), Inches(12.23), RULE)
hline(s, Inches(0.55), Inches(3.35), Inches(12.23), RULE_2)

stats = [
    ("9", "SPECIALISED AI AGENTS", True),
    ("30+", "CLAUSE TYPES EXTRACTED", True),
    ("10", "RISK DIMENSIONS SCORED", False),
    ("13", "SLA METRIC TYPES", True),
]
stat_w = Inches(3.0)
stat_x0 = Inches(0.55)
for i, (v, k, accent) in enumerate(stats):
    x = stat_x0 + i * stat_w
    if i > 0:
        vline(s, x, Inches(2.62), Inches(0.6), RULE_2)
    tx(s, x + Inches(0.18), Inches(2.62), Inches(2.5), Inches(0.5),
       v, font=FONT_SANS, size=30, color=(ORANGE_2 if accent else INK), bold=True)
    tx(s, x + Inches(0.18), Inches(3.08), Inches(2.6), Inches(0.25),
       k, font=FONT_MONO, size=9, color=INK_3, spacing=1.6)

# Capabilities 3×2 grid
caps = [
    ("01", "Contract Intelligence",   "Metadata, 30+ clause types, 10 risk categories, semantic search.", False),
    ("02", "Obligation Tracking",     "RAG status, owner assignment, escalation, evidence.", False),
    ("03", "SLA Monitoring",          "Breach detection, automated credit math, vendor scorecards.", False),
    ("04", "Renewal Management",      "Auto-renewal detection, notice tracking, AI recommendations.", False),
    ("05", "Relationship Governance", "Perception gap analysis, composite health, improvement points.", True),
    ("06", "Enterprise Security",     "Encryption, RBAC, multi-tenant isolation, audit trails.", False),
]
cap_w = Inches(4.05); cap_h = Inches(1.55)
cap_x0 = Inches(0.55); cap_y0 = Inches(3.65)
cap_gap = Inches(0.05)

for i, (n, t, d, feat) in enumerate(caps):
    row = i // 3; col = i % 3
    x = cap_x0 + col * (cap_w + cap_gap)
    y = cap_y0 + row * (cap_h + cap_gap)

    if feat:
        rect(s, x, y, cap_w, cap_h, ORANGE_2, line_color=ORANGE_2, line_width=Pt(0.75))
        n_color = RGBColor(0xF7, 0xF6, 0xF3); t_color = PAPER; d_color = RGBColor(0xED, 0xE5, 0xD8)
    else:
        rect(s, x, y, cap_w, cap_h, SURFACE, line_color=RULE_2, line_width=Pt(0.75))
        n_color = ORANGE_2; t_color = INK; d_color = INK_3

    tx(s, x + Inches(0.22), y + Inches(0.18), Inches(2), Inches(0.2),
       n, font=FONT_MONO, size=9, color=n_color, spacing=1.6, bold=True)
    tx(s, x + Inches(0.22), y + Inches(0.42), cap_w - Inches(0.45), Inches(0.35),
       t, font=FONT_SANS, size=15, color=t_color, bold=True)
    tx(s, x + Inches(0.22), y + Inches(0.78), cap_w - Inches(0.45), Inches(0.7),
       d, font=FONT_SANS, size=11, color=d_color, line_spacing=1.3)
    if feat:
        # "Only in Evaluetor" tag
        tag_y = y + Inches(1.2)
        rect(s, x + Inches(0.22), tag_y, Inches(1.5), Inches(0.22), ORANGE_2, line_color=PAPER, line_width=Pt(0.75))
        tx(s, x + Inches(0.22), tag_y + Inches(0.01), Inches(1.5), Inches(0.22),
           "ONLY IN EVALUETOR", font=FONT_MONO, size=7.5, color=PAPER, spacing=2, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 03 · THE FOUR DISCIPLINES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "03 / THE FOUR DISCIPLINES", "Slide 03 — The four disciplines", bg=PAPER_2)

# Headline
tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Four disciplines, ", {"size": 32, "color": INK, "font": FONT_SANS}),
    ("examined in turn.", {"size": 32, "color": INK_3, "font": FONT_ITAL, "italic": True})])

tx(s, Inches(0.55), Inches(1.78), Inches(12), Inches(0.4),
   "What each one does, what it extracts, and what it produces — built in, not bolted on.",
   font=FONT_SANS, size=13, color=INK_3)

# 2×2 grid of disciplines
disciplines = [
    ("01", "Obligation Tracking", False,
     "Every promise — tracked, owned, evidenced.",
     ["Classification across 7 types (payment, delivery, reporting, notification, performance, compliance, other)",
      "RAG status dashboard, owner assignment, escalation paths",
      "Configurable deadline alerts to email · Slack · Teams",
      "Evidence attachments, compliance reporting, audit-ready trail"]),
    ("02", "Risk Detection", False,
     "What is buried in the boilerplate.",
     ["10 risk patterns: uncapped liability, broad indemnification, weak termination, auto-renewal traps, unfavorable IP, weak confidentiality, missing liability caps, one-sided terms, regulatory exposure, ambiguous language",
      "4-level severity scoring (Low / Med / High / Critical)",
      "Knowledge-graph anomaly detection vs portfolio norms",
      "Portfolio risk dashboard with drill-down by vendor and type"]),
    ("03", "SLA Monitoring", False,
     "Recover what you are owed — automatically.",
     ["13 SLA metric types: uptime, response, resolution, throughput…",
      "Real-time breach detection across the vendor base",
      "Automated service-credit calculation per contract terms",
      "A-to-F vendor scorecards: compliance · breaches · impact"]),
    ("04", "Relationship Governance", True,
     "The view from both sides of the table.",
     ["Perception-gap analysis: how you see them vs they see you",
      "Composite health: Risk 30% · SLA 40% · Obligations 30%",
      "Auto-generated improvement points with owners and priorities",
      "Multi-party satisfaction surveys, KPI tracking by relationship"]),
]

card_w = Inches(6.1); card_h = Inches(2.3)
card_x0 = Inches(0.55); card_y0 = Inches(2.42)
card_gap = Inches(0.13)

for i, (n, t, exclusive, tag, bullets) in enumerate(disciplines):
    row = i // 2; col = i % 2
    x = card_x0 + col * (card_w + card_gap)
    y = card_y0 + row * (card_h + card_gap)
    rect(s, x, y, card_w, card_h, SURFACE, line_color=RULE_2, line_width=Pt(0.75))

    # Big numeral
    tx(s, x + Inches(0.3), y + Inches(0.18), Inches(1), Inches(0.6),
       n, font=FONT_SANS, size=36, color=ORANGE_2, bold=False)
    # Discipline name
    name_runs = [(t, {"size": 17, "color": INK, "font": FONT_SANS, "bold": True})]
    if exclusive:
        name_runs.append(("  EXCLUSIVE", {"size": 8, "color": ORANGE_2, "font": FONT_MONO, "bold": True}))
    tx(s, x + Inches(1.3), y + Inches(0.22), card_w - Inches(1.5), Inches(0.35),
       name_runs)
    # Italic tag
    tx(s, x + Inches(1.3), y + Inches(0.55), card_w - Inches(1.5), Inches(0.3),
       tag, font=FONT_ITAL, size=12, color=INK_2, italic=True)

    # Hairline
    hline(s, x + Inches(0.3), y + Inches(0.96), card_w - Inches(0.6), RULE_2)

    # Bullets
    for j, b in enumerate(bullets):
        by = y + Inches(1.04) + j * Inches(0.27)
        tx(s, x + Inches(0.3), by, Inches(0.2), Inches(0.25),
           "→", font=FONT_MONO, size=10, color=INK_3)
        tx(s, x + Inches(0.55), by, card_w - Inches(0.8), Inches(0.3),
           b, font=FONT_SANS, size=11, color=INK_2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 04 · BUILT FOR TRUST
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "04 / BUILT FOR TRUST", "Slide 04 — Built for trust", bg=PAPER)

# Headline
tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Connected. Inspectable. ", {"size": 32, "color": INK, "font": FONT_SANS}),
    ("Built for European enterprise.", {"size": 32, "color": INK_3, "font": FONT_ITAL, "italic": True})])

tx(s, Inches(0.55), Inches(1.78), Inches(12), Inches(0.4),
   "Reads operational reality. Traces every AI decision. Designed multi-tenant from day one.",
   font=FONT_SANS, size=13, color=INK_3)

# 3 columns
trust_cols = [
    ("INTEGRATION SCOPE", "Read the world",
     "Compare contracted terms against the systems where the work actually happens.",
     [("ServiceNow ITSM", "Shipping"),
      ("ERP volumes (SAP · Oracle)", "Per engagement"),
      ("FX rates (ECB feeds)", "Per engagement"),
      ("SharePoint document store", "Shipping"),
      ("REST API + webhooks", "All endpoints")]),
    ("AI-NATIVE ARCHITECTURE", "Inspectable AI",
     "Nine agents in parallel, traced and replayable. Defensible for risk teams.",
     [("Specialised AI agents", "9 native"),
      ("Knowledge graph anomalies", "Built-in"),
      ("Langfuse observability", "Per call"),
      ("Human-in-the-loop", "First 30 days"),
      ("Source-clause citation", "Every flag")]),
    ("ENTERPRISE SECURITY", "European by default",
     "Multi-tenant isolation at the data layer. GDPR-aware. EU residency available.",
     [("Multi-tenant isolation", "Data layer"),
      ("Role-based access control", "API layer"),
      ("Encryption at rest & transit", "Standard"),
      ("Full audit trail", "Every action"),
      ("EU residency", "Available"),
      ("SOC 2 · ISO 27001", "Roadmap")]),
]

col_w = Inches(4.05); col_h = Inches(4.5)
col_x0 = Inches(0.55); col_y0 = Inches(2.4)
col_gap = Inches(0.05)

for i, (label, title, ctx, rows) in enumerate(trust_cols):
    x = col_x0 + i * (col_w + col_gap)
    rect(s, x, col_y0, col_w, col_h, SURFACE, line_color=RULE_2, line_width=Pt(0.75))

    tx(s, x + Inches(0.25), col_y0 + Inches(0.22), col_w - Inches(0.5), Inches(0.25),
       label, font=FONT_MONO, size=9, color=ORANGE_2, spacing=1.6, bold=True)
    tx(s, x + Inches(0.25), col_y0 + Inches(0.52), col_w - Inches(0.5), Inches(0.35),
       title, font=FONT_SANS, size=17, color=INK, bold=True)
    tx(s, x + Inches(0.25), col_y0 + Inches(0.88), col_w - Inches(0.5), Inches(0.6),
       ctx, font=FONT_ITAL, size=11.5, color=INK_2, italic=True, line_spacing=1.3)
    hline(s, x + Inches(0.25), col_y0 + Inches(1.6), col_w - Inches(0.5), RULE_2)

    for j, (k, v) in enumerate(rows):
        ry = col_y0 + Inches(1.75) + j * Inches(0.42)
        tx(s, x + Inches(0.25), ry, col_w - Inches(1.5), Inches(0.3),
           k, font=FONT_SANS, size=11, color=INK_2)
        tx(s, x + col_w - Inches(1.5), ry, Inches(1.25), Inches(0.3),
           v, font=FONT_MONO, size=10, color=ORANGE_2, bold=True, align=PP_ALIGN.RIGHT, spacing=0.5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 05 · WHAT THE OTHERS DO NOT DO (comparison table)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "05 / COMPARED", "Slide 05 — What the others do not do", bg=PAPER_2)

# Headline
tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("What the legacy CLM tools ", {"size": 32, "color": INK, "font": FONT_SANS}),
    ("do not do.", {"size": 32, "color": INK_3, "font": FONT_ITAL, "italic": True})])

tx(s, Inches(0.55), Inches(1.78), Inches(12), Inches(0.4),
   "Evaluetor was built model-first. The difference is visible in the cells of this table.",
   font=FONT_SANS, size=13, color=INK_3)

# Comparison table
headers = ["Capability", "Evaluetor", "DocuSign CLM", "Icertis", "Ironclad"]
rows = [
    ("Native AI agents",         "Nine",          "Limited",      "Add-on",       "Limited"),
    ("SLA monitoring",           "Built-in",      "Not standard", "Add-on",       "Not standard"),
    ("Obligation tracking",      "AI-extracted",  "Manual",       "Manual",       "Limited"),
    ("Relationship governance",  "Built-in",      "Not standard", "Not standard", "Not standard"),
    ("Knowledge graph",          "Built-in",      "Not standard", "Not standard", "Not standard"),
    ("LLM observability",        "Langfuse",      "Not standard", "Not standard", "Not standard"),
    ("Unlimited users",          "Yes",           "Per-seat",     "Per-seat",     "Per-seat"),
    ("Time to first outcome",    "Within weeks",  "6–12 months",  "12–18 months", "3–6 months"),
]

table_x0 = Inches(0.55); table_y0 = Inches(2.5)
col_widths = [Inches(3.3), Inches(2.4), Inches(2.2), Inches(2.0), Inches(2.0)]

# Background for the table
table_w = sum(col_widths, Emu(0))
table_h = Inches(0.55) + Inches(0.46) * len(rows)
rect(s, table_x0, table_y0, table_w, table_h, SURFACE, line_color=RULE_2, line_width=Pt(0.75))

# Header row
cx = table_x0
for i, h in enumerate(headers):
    if i == 1:
        # Evaluetor: orange dot + label
        rect(s, cx + Inches(0.18), table_y0 + Inches(0.21), Inches(0.08), Inches(0.08), ORANGE_2)
        tx(s, cx + Inches(0.32), table_y0 + Inches(0.16), col_widths[i] - Inches(0.4), Inches(0.3),
           h, font=FONT_SANS, size=12, color=INK, bold=True)
    elif i == 0:
        tx(s, cx + Inches(0.18), table_y0 + Inches(0.18), col_widths[i] - Inches(0.4), Inches(0.3),
           h.upper(), font=FONT_MONO, size=8.5, color=INK_3, spacing=1.4, bold=True)
    else:
        tx(s, cx + Inches(0.18), table_y0 + Inches(0.18), col_widths[i] - Inches(0.4), Inches(0.3),
           h.upper(), font=FONT_MONO, size=8.5, color=INK_3, spacing=1.4, bold=True)
    cx += col_widths[i]
hline(s, table_x0, table_y0 + Inches(0.55), table_w, RULE)

# Body rows
ry = table_y0 + Inches(0.55)
for ri, row in enumerate(rows):
    if ri % 2 == 1:
        rect(s, table_x0, ry, table_w, Inches(0.46), PAPER_2, line_color=PAPER_2, line_width=Pt(0))
    cx = table_x0
    for ci, cell in enumerate(row):
        if ci == 0:
            tx(s, cx + Inches(0.18), ry + Inches(0.12), col_widths[ci] - Inches(0.4), Inches(0.3),
               cell, font=FONT_SANS, size=12, color=INK, bold=True)
        elif ci == 1:
            tx(s, cx + Inches(0.18), ry + Inches(0.12), col_widths[ci] - Inches(0.4), Inches(0.3),
               cell, font=FONT_SANS, size=12, color=INK, bold=True)
        else:
            tx(s, cx + Inches(0.18), ry + Inches(0.12), col_widths[ci] - Inches(0.4), Inches(0.3),
               cell, font=FONT_SANS, size=12, color=INK_3)
        cx += col_widths[ci]
    ry += Inches(0.46)
    if ri < len(rows) - 1:
        hline(s, table_x0, ry, table_w, RULE_2, thickness=Pt(0.5))

# Note
tx(s, Inches(0.55), Inches(6.6), Inches(12), Inches(0.3),
   "Comparison reflects publicly documented capabilities at time of writing (May 2026). "
   "Vendor packaging and tiering vary; \"Not standard\" indicates not part of the default offering.",
   font=FONT_MONO, size=8.5, color=INK_3, line_spacing=1.5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 06 · THANK YOU
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "CLOSE", "Slide 06 — Close", bg=PAPER_2)

# Big orange E mark
rect(s, Inches(0.7), Inches(1.4), Inches(0.85), Inches(0.85), ORANGE_2)
tx(s, Inches(0.7), Inches(1.5), Inches(0.85), Inches(0.7),
   "E", font=FONT_MONO, size=32, color=PAPER, bold=True, align=PP_ALIGN.CENTER)

# "Thank you." headline
tx(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.2),
   [("Thank you. ", {"size": 72, "color": INK, "font": FONT_SANS}),
    ("For your time.", {"size": 72, "color": INK_3, "font": FONT_ITAL, "italic": True})])

# Invitation paragraph
tx(s, Inches(0.7), Inches(4.0), Inches(11.5), Inches(1),
   "Bring five contracts. We will show you, in thirty minutes, what your current CLM is missing. "
   "No payment, no procurement, no preparation required — findings are yours to keep.",
   font=FONT_SANS, size=17, color=INK_2, line_spacing=1.45)

# Contact bar — 3 columns
hline(s, Inches(0.7), Inches(5.5), Inches(11.5), RULE_2)

# Steven Visser
tx(s, Inches(0.7), Inches(5.7), Inches(3), Inches(0.25),
   "STEVEN VISSER", font=FONT_MONO, size=9, color=ORANGE_2, spacing=1.6, bold=True)
tx(s, Inches(0.7), Inches(6.0), Inches(3), Inches(0.3),
   "Managing Partner", font=FONT_SANS, size=13, color=INK, bold=True)

# Email
tx(s, Inches(4.5), Inches(5.7), Inches(4), Inches(0.25),
   "EMAIL", font=FONT_MONO, size=9, color=ORANGE_2, spacing=1.6, bold=True)
tx(s, Inches(4.5), Inches(6.0), Inches(5), Inches(0.3),
   "Steven.Visser@Evaluetor.com", font=FONT_MONO, size=12, color=INK, bold=True)

# Voice
tx(s, Inches(9.3), Inches(5.7), Inches(3), Inches(0.25),
   "VOICE", font=FONT_MONO, size=9, color=ORANGE_2, spacing=1.6, bold=True)
tx(s, Inches(9.3), Inches(6.0), Inches(3.5), Inches(0.3),
   "+31 6 31 00 34 95", font=FONT_MONO, size=12, color=INK, bold=True)


# ─── Save ────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Evaluetor-Capabilities-Briefing.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
