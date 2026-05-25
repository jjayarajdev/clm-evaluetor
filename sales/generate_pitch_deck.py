"""Generate Evaluetor Sales Pitch Deck — Dark (with orange) + Light versions."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ───────────────────────────────────────────────────────────────
# COLOR THEMES
# ───────────────────────────────────────────────────────────────

DARK_THEME = dict(
    bg        = RGBColor(0x03, 0x07, 0x12),
    card      = RGBColor(0x11, 0x18, 0x27),
    heading   = RGBColor(0xFF, 0xFF, 0xFF),
    body      = RGBColor(0x9C, 0xA3, 0xAF),
    muted     = RGBColor(0x6B, 0x72, 0x80),
    label     = RGBColor(0xD1, 0xD5, 0xDB),
    border    = RGBColor(0x2D, 0x33, 0x48),
    connector = RGBColor(0x37, 0x3D, 0x55),
    tbl_hdr   = RGBColor(0x0A, 0x0F, 0x1A),
    tbl_alt   = RGBColor(0x11, 0x18, 0x27),
    circle_bg = RGBColor(0x03, 0x07, 0x12),
    step_num  = RGBColor(0x03, 0x07, 0x12),
    # accents
    purple    = RGBColor(0x7C, 0x3A, 0xED),
    purple_lt = RGBColor(0xA7, 0x8B, 0xFA),
    orange    = RGBColor(0xFB, 0x92, 0x3C),
    orange_lt = RGBColor(0xFD, 0xBA, 0x74),
    pink      = RGBColor(0xEC, 0x48, 0x99),
    red       = RGBColor(0xF8, 0x71, 0x71),
    emerald   = RGBColor(0x34, 0xD3, 0x99),
    blue      = RGBColor(0x60, 0xA5, 0xFA),
    amber     = RGBColor(0xFB, 0xBF, 0x24),
    cyan      = RGBColor(0x22, 0xD3, 0xEE),
)

LIGHT_THEME = dict(
    bg        = RGBColor(0xFF, 0xFF, 0xFF),
    card      = RGBColor(0xF9, 0xFA, 0xFB),
    heading   = RGBColor(0x11, 0x18, 0x27),
    body      = RGBColor(0x4B, 0x55, 0x63),
    muted     = RGBColor(0x6B, 0x72, 0x80),
    label     = RGBColor(0x37, 0x41, 0x51),
    border    = RGBColor(0xE5, 0xE7, 0xEB),
    connector = RGBColor(0xD1, 0xD5, 0xDB),
    tbl_hdr   = RGBColor(0xF3, 0xF4, 0xF6),
    tbl_alt   = RGBColor(0xF9, 0xFA, 0xFB),
    circle_bg = RGBColor(0xFF, 0xFF, 0xFF),
    step_num  = RGBColor(0xFF, 0xFF, 0xFF),
    # accents — slightly deeper for light-bg contrast
    purple    = RGBColor(0x7C, 0x3A, 0xED),
    purple_lt = RGBColor(0x6D, 0x28, 0xD9),
    orange    = RGBColor(0xEA, 0x76, 0x0D),
    orange_lt = RGBColor(0xF9, 0x73, 0x16),
    pink      = RGBColor(0xDB, 0x27, 0x77),
    red       = RGBColor(0xDC, 0x26, 0x26),
    emerald   = RGBColor(0x05, 0x96, 0x69),
    blue      = RGBColor(0x25, 0x63, 0xEB),
    amber     = RGBColor(0xD9, 0x77, 0x06),
    cyan      = RGBColor(0x08, 0x91, 0xB2),
)

TOTAL_SLIDES = 12
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ───────────────────────────────────────────────────────────────
# HELPERS (theme-aware)
# ───────────────────────────────────────────────────────────────

def set_slide_bg(slide, t):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = t["bg"]


def add_shape_fill(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def add_rounded_rect(slide, left, top, width, height, t, fill_key="card"):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.line.color.rgb = t["border"]
    s.line.width = Pt(1)
    s.fill.solid()
    s.fill.fore_color.rgb = t[fill_key]
    return s


def add_text(slide, left, top, width, height, text, font_size=18,
             color=None, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_tag(slide, left, top, text, color):
    return add_text(slide, left, top, Inches(4), Inches(0.35),
                    text.upper(), font_size=11, color=color, bold=True)


def add_slide_number(slide, num, t):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num} / {TOTAL_SLIDES}", font_size=10, color=t["muted"],
             alignment=PP_ALIGN.RIGHT)


def add_dot(slide, x, y, color, size=Inches(0.12)):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    d.line.fill.background()
    d.fill.solid()
    d.fill.fore_color.rgb = color


def add_connector(slide, x, y, height, t):
    c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.04), height)
    c.line.fill.background()
    c.fill.solid()
    c.fill.fore_color.rgb = t["connector"]


def accent_strip(slide, t):
    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), t["purple"])


def accent_strip_orange(slide, t):
    """Orange-tinted accent strip — half purple, half orange."""
    add_shape_fill(slide, Inches(0), Inches(0), Inches(6.667), Inches(0.05), t["purple"])
    add_shape_fill(slide, Inches(6.667), Inches(0), Inches(6.666), Inches(0.05), t["orange"])


# ───────────────────────────────────────────────────────────────
# BUILD THE FULL DECK
# ───────────────────────────────────────────────────────────────

def build_deck(t, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── Deep-dive builder (closure over prs, t) ──
    def build_deep_dive(slide_num, tag_text, tag_color, title, subtitle,
                        left_title, left_subtitle, pipeline_steps,
                        right_title, feature_cards, use_orange_strip=False):
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide, t)
        (accent_strip_orange if use_orange_strip else accent_strip)(slide, t)

        add_tag(slide, Inches(0.8), Inches(0.6), tag_text, tag_color)
        add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                 title, font_size=40, color=t["heading"], bold=True)
        add_text(slide, Inches(0.8), Inches(1.8), Inches(9), Inches(0.7),
                 subtitle, font_size=16, color=t["body"])

        add_text(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(0.4),
                 left_title, font_size=18, color=t["heading"], bold=True)
        add_text(slide, Inches(0.8), Inches(3.35), Inches(5.5), Inches(0.5),
                 left_subtitle, font_size=13, color=t["body"])

        for i, (st, sd, ac) in enumerate(pipeline_steps):
            y = Inches(4.0) + Inches(i * 0.65)
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                             Inches(1.0), y + Inches(0.05),
                                             Inches(0.3), Inches(0.3))
            circle.line.fill.background()
            circle.fill.solid()
            circle.fill.fore_color.rgb = ac
            add_text(slide, Inches(1.0), y + Inches(0.06), Inches(0.3), Inches(0.3),
                     str(i + 1), font_size=11, color=t["step_num"], bold=True,
                     alignment=PP_ALIGN.CENTER)
            add_text(slide, Inches(1.5), y, Inches(2.0), Inches(0.3),
                     st, font_size=13, color=t["heading"], bold=True)
            add_text(slide, Inches(3.5), y + Inches(0.02), Inches(3.0), Inches(0.3),
                     sd, font_size=11, color=t["body"])
            if i < len(pipeline_steps) - 1:
                add_connector(slide, Inches(1.13), y + Inches(0.38), Inches(0.28), t)

        rx = Inches(7.2)
        add_text(slide, rx, Inches(2.9), Inches(5.5), Inches(0.4),
                 right_title, font_size=18, color=t["heading"], bold=True)
        for i, (ft, fd, ac) in enumerate(feature_cards):
            y = Inches(3.5) + Inches(i * 1.35)
            add_rounded_rect(slide, rx, y, Inches(5.3), Inches(1.2), t)
            add_dot(slide, rx + Inches(0.25), y + Inches(0.25), ac)
            add_text(slide, rx + Inches(0.5), y + Inches(0.15),
                     Inches(4.5), Inches(0.3), ft,
                     font_size=14, color=t["heading"], bold=True)
            add_text(slide, rx + Inches(0.25), y + Inches(0.5),
                     Inches(4.8), Inches(0.7), fd,
                     font_size=11, color=t["body"])

        add_slide_number(slide, slide_num, t)
        return slide

    # ================================================================
    # SLIDE 1: TITLE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip_orange(s, t)

    badge = add_rounded_rect(s, Inches(5.5), Inches(1.6), Inches(2.3), Inches(0.5), t)
    add_text(s, Inches(5.5), Inches(1.62), Inches(2.3), Inches(0.5),
             "EVALUETOR", font_size=14, color=t["purple_lt"], bold=True,
             alignment=PP_ALIGN.CENTER)

    add_text(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
             "Your contracts hold the answers.", font_size=44,
             color=t["heading"], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.2),
             "Now you can find them.", font_size=44,
             color=t["orange"], bold=True, alignment=PP_ALIGN.CENTER)

    add_text(s, Inches(2.5), Inches(4.7), Inches(8.3), Inches(0.9),
             "AI-native contract intelligence that makes your contracts transparent, "
             "measurable, and actionable \u2014 in minutes, not months.",
             font_size=18, color=t["body"], alignment=PP_ALIGN.CENTER)

    pillars = [
        ("Contract Intelligence", t["purple_lt"]),
        ("Relationship Governance", t["pink"]),
        ("Enterprise Security", t["orange"]),
    ]
    for i, (label, clr) in enumerate(pillars):
        x = Inches(3.5) + Inches(i * 2.3)
        add_dot(s, x, Inches(5.95), clr)
        add_text(s, x + Inches(0.2), Inches(5.85), Inches(2.0), Inches(0.35),
                 label, font_size=11, color=t["muted"])

    add_text(s, Inches(0.8), Inches(6.9), Inches(4), Inches(0.4),
             "evaluetor.com  |  AI-Powered Contract Intelligence",
             font_size=10, color=t["muted"])
    add_slide_number(s, 1, t)

    # ================================================================
    # SLIDE 2: THE CHALLENGE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip(s, t)

    add_tag(s, Inches(0.8), Inches(0.6), "The Challenge", t["red"])
    add_text(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "The gap between signing and managing", font_size=40,
             color=t["heading"], bold=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.7),
             "Contracts are negotiated carefully, then stored away. The obligations, "
             "SLAs, and renewal terms inside them often go unmonitored \u2014 "
             "leaving real value on the table.",
             font_size=16, color=t["body"])

    problems = [
        ("60%", "Limited Visibility",
         "of enterprises struggle to surface\nwhat\u2019s actually inside their\ncontract portfolio.",
         t["orange"]),
        ("$2M+", "Value at Risk",
         "lost annually to missed obligations,\nunnoticed auto-renewals, and\npenalties that could be avoided.",
         t["red"]),
        ("15\u201330%", "Unclaimed Credits",
         "of SLA service credits go unrecovered.\nBreaches happen \u2014 the challenge\nis detecting them in time.",
         t["orange"]),
    ]
    for i, (stat, title, desc, stat_clr) in enumerate(problems):
        x = Inches(0.8) + Inches(i * 4.0)
        add_rounded_rect(s, x, Inches(3.0), Inches(3.6), Inches(3.5), t)
        add_text(s, x + Inches(0.4), Inches(3.3), Inches(2.8), Inches(0.8),
                 stat, font_size=48, color=stat_clr, bold=True)
        add_text(s, x + Inches(0.4), Inches(4.2), Inches(2.8), Inches(0.5),
                 title, font_size=20, color=t["heading"], bold=True)
        add_text(s, x + Inches(0.4), Inches(4.8), Inches(2.8), Inches(1.2),
                 desc, font_size=14, color=t["body"])

    add_slide_number(s, 2, t)

    # ================================================================
    # SLIDE 3: THE SOLUTION
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip_orange(s, t)

    add_tag(s, Inches(0.8), Inches(0.6), "The Solution", t["purple_lt"])
    add_text(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "One platform. Complete intelligence.", font_size=40,
             color=t["heading"], bold=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.5),
             "9 specialized AI agents extract, analyze, and monitor everything "
             "in your contracts \u2014 automatically.",
             font_size=16, color=t["body"])

    capabilities = [
        ("Contract Intelligence", "AI extracts metadata, 30+ clause types, "
         "and risks across 10 categories. Full semantic search across entire portfolio.",
         t["purple_lt"]),
        ("Obligation Tracking", "RAG status indicators, owner assignment, "
         "escalation paths, and automated deadline alerts. Never miss a commitment.",
         t["emerald"]),
        ("SLA Monitoring", "Real-time breach detection, service credit "
         "calculation, and vendor scorecards. Know when performance falls short.",
         t["blue"]),
        ("Renewal Management", "Auto-renewal detection, notice period "
         "tracking, and AI-powered renewal recommendations. Stay ahead of every date.",
         t["orange"]),
        ("Relationship Governance", "KPI perception gap analysis, "
         "vendor health scores, satisfaction surveys, and improvement tracking. "
         "Only available in Evaluetor.",
         t["pink"]),
        ("Enterprise Security", "SOC 2-ready architecture, GDPR-aware, AES-256 "
         "encryption, role-based access, multi-tenant isolation, full audit trails.",
         t["cyan"]),
    ]

    for i, (title, desc, accent) in enumerate(capabilities):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + Inches(col * 4.0)
        y = Inches(2.7) + Inches(row * 2.4)
        add_rounded_rect(s, x, y, Inches(3.6), Inches(2.1), t)
        add_dot(s, x + Inches(0.3), y + Inches(0.35), accent, Inches(0.14))
        add_text(s, x + Inches(0.55), y + Inches(0.25),
                 Inches(2.8), Inches(0.35), title,
                 font_size=16, color=t["heading"], bold=True)
        add_text(s, x + Inches(0.3), y + Inches(0.75),
                 Inches(3.0), Inches(1.2), desc,
                 font_size=12, color=t["body"])

    add_slide_number(s, 3, t)

    # ================================================================
    # SLIDE 4: HOW IT WORKS
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip(s, t)

    add_tag(s, Inches(0.8), Inches(0.6), "How It Works", t["purple_lt"])
    add_text(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "From upload to intelligence in minutes", font_size=40,
             color=t["heading"], bold=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.5),
             "No templates. No pre-configuration. No training required. Just upload and go.",
             font_size=16, color=t["body"])

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(2.2), Inches(4.05), Inches(9.0), Inches(0.03))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = t["connector"]

    steps = [
        ("01", "Upload", "Drop any PDF, DOCX, or\nbatch upload a ZIP.\nOCR handles scans.", t["purple_lt"]),
        ("02", "AI Extracts", "9 agents extract metadata,\nclauses, obligations, SLAs,\nrisks, and renewals.", t["orange"]),
        ("03", "Dashboards", "Real-time dashboards surface\nrisks, deadlines, SLA breaches,\nand vendor scores.", t["pink"]),
        ("04", "Take Action", "Claim credits, meet deadlines,\nrenegotiate terms, govern\nrelationships with data.", t["emerald"]),
    ]

    for i, (num, title, desc, accent) in enumerate(steps):
        x = Inches(1.0) + Inches(i * 3.0)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.7), Inches(3.3),
                                     Inches(1.1), Inches(1.1))
        circle.line.color.rgb = accent
        circle.line.width = Pt(2)
        circle.fill.solid()
        circle.fill.fore_color.rgb = t["circle_bg"]
        add_text(s, x + Inches(0.7), Inches(3.45), Inches(1.1), Inches(0.9),
                 num, font_size=28, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(4.7), Inches(2.1), Inches(0.4),
                 title, font_size=18, color=t["heading"], bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.05), Inches(5.2), Inches(2.4), Inches(1.2),
                 desc, font_size=12, color=t["body"], alignment=PP_ALIGN.CENTER)

    add_slide_number(s, 4, t)

    # ================================================================
    # SLIDE 5: UNDER THE HOOD — ARCHITECTURE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip_orange(s, t)

    add_tag(s, Inches(0.5), Inches(0.25), "Architecture", t["orange"])
    add_text(s, Inches(0.5), Inches(0.55), Inches(12), Inches(0.5),
             "Under the hood: what happens when you upload a contract",
             font_size=30, color=t["heading"], bold=True)

    # ── ROW 1: AI PIPELINE (left to right) ──
    add_text(s, Inches(0.5), Inches(1.2), Inches(3), Inches(0.3),
             "STAGE 1: AI Pipeline", font_size=13, color=t["purple_lt"], bold=True)
    add_text(s, Inches(3.5), Inches(1.23), Inches(6), Inches(0.3),
             "9 specialized agents process every contract automatically",
             font_size=10, color=t["body"])

    agents = [
        ("Upload\n& Parse", t["purple_lt"]),
        ("Metadata\nExtraction", t["blue"]),
        ("Risk\nDetection", t["red"]),
        ("Clause\nExtraction", t["purple_lt"]),
        ("Obligation\nExtraction", t["emerald"]),
        ("SLA\nExtraction", t["blue"]),
        ("Renewal\nDetection", t["orange"]),
        ("Auto-Link\nDetection", t["pink"]),
        ("Q&A\nReady", t["emerald"]),
    ]
    for i, (label, accent) in enumerate(agents):
        x = Inches(0.4) + Inches(i * 1.4)
        card = add_rounded_rect(s, x, Inches(1.55), Inches(1.25), Inches(0.85), t)
        card.line.color.rgb = accent
        add_text(s, x + Inches(0.05), Inches(1.6), Inches(1.15), Inches(0.75),
                 label, font_size=10, color=t["heading"], bold=True,
                 alignment=PP_ALIGN.CENTER)
        # Arrow between agents
        if i < len(agents) - 1:
            add_text(s, x + Inches(1.22), Inches(1.8), Inches(0.2), Inches(0.3),
                     "\u2192", font_size=11, color=t["connector"], bold=True,
                     alignment=PP_ALIGN.CENTER)

    # ── ROW 2: GOLDEN DATASET (feedback loop) ──
    add_text(s, Inches(0.5), Inches(2.65), Inches(3), Inches(0.3),
             "STAGE 2: Golden Dataset", font_size=13, color=t["amber"], bold=True)
    add_text(s, Inches(3.5), Inches(2.68), Inches(6), Inches(0.3),
             "Verified extractions improve accuracy over time",
             font_size=10, color=t["body"])

    golden_cards = [
        ("Verified Contracts",
         "50+ representative contracts\nwith human-verified extractions",
         t["amber"]),
        ("Few-Shot Learning",
         "Verified examples injected into\nAI prompts for higher accuracy",
         t["orange"]),
        ("Quality Metrics",
         "Precision, recall, F1 per field\ntracked across contract types",
         t["emerald"]),
        ("Continuous Improvement",
         "Each correction feeds back \u2014\nthe system gets smarter over time",
         t["purple_lt"]),
    ]
    for i, (title, desc, accent) in enumerate(golden_cards):
        x = Inches(0.4) + Inches(i * 3.15)
        card = add_rounded_rect(s, x, Inches(3.0), Inches(2.95), Inches(1.0), t)
        card.line.color.rgb = accent
        add_dot(s, x + Inches(0.12), Inches(3.12), accent, Inches(0.1))
        add_text(s, x + Inches(0.3), Inches(3.05), Inches(2.5), Inches(0.25),
                 title, font_size=11, color=t["heading"], bold=True)
        add_text(s, x + Inches(0.12), Inches(3.35), Inches(2.7), Inches(0.6),
                 desc, font_size=9, color=t["body"])

    # Feedback arrow
    add_text(s, Inches(10.5), Inches(2.65), Inches(2.5), Inches(0.3),
             "\u21bb feeds back into Stage 1", font_size=9, color=t["amber"],
             bold=True)

    # ── ROW 3: GOVERNANCE BRIDGE (the USP) ──
    add_text(s, Inches(0.5), Inches(4.25), Inches(3.5), Inches(0.3),
             "STAGE 3: Governance Bridge", font_size=13, color=t["pink"], bold=True)
    add_text(s, Inches(4.0), Inches(4.28), Inches(6), Inches(0.3),
             "Unique to Evaluetor \u2014 no other CLM does this",
             font_size=10, color=t["pink"], bold=True)

    gov_steps = [
        ("Counterparty\n\u2192 Organization", "Auto-create or\nmatch existing org", t["blue"]),
        ("Contract\n\u2192 Relationship", "Business relationship\nwith health scoring", t["emerald"]),
        ("SLA Metrics\n\u2192 KPIs", "Targets & thresholds\nauto-created", t["orange"]),
        ("Risk Clauses\n\u2192 Improvements", "Action items from\nhigh-risk clauses", t["red"]),
        ("Health Score\nCalculation", "Risk 30% + SLA 40%\n+ Obligations 30%", t["pink"]),
        ("Perception\nGap Scoring", "Internal vs external\nview of performance", t["purple_lt"]),
    ]
    for i, (title, desc, accent) in enumerate(gov_steps):
        x = Inches(0.4) + Inches(i * 2.1)
        card = add_rounded_rect(s, x, Inches(4.65), Inches(1.95), Inches(1.2), t)
        card.line.color.rgb = accent
        add_text(s, x + Inches(0.08), Inches(4.7), Inches(1.8), Inches(0.5),
                 title, font_size=10, color=accent, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.08), Inches(5.2), Inches(1.8), Inches(0.55),
                 desc, font_size=9, color=t["body"], alignment=PP_ALIGN.CENTER)

    # ── BOTTOM: Query bar ──
    qbar = add_rounded_rect(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.65), t)
    qbar.line.color.rgb = t["emerald"]
    add_text(s, Inches(0.7), Inches(6.15), Inches(1.5), Inches(0.55),
             "Ask Anything", font_size=14, color=t["emerald"], bold=True)
    add_text(s, Inches(2.3), Inches(6.18), Inches(10.3), Inches(0.5),
             "\u201cWhat obligations are due next month?\u201d  \u2022  "
             "\u201cShow me all high-risk indemnification clauses\u201d  \u2022  "
             "\u201cWhich vendors have SLA breaches?\u201d  \u2022  "
             "\u201cCompare payment terms across all MSAs\u201d",
             font_size=10, color=t["body"])

    add_slide_number(s, 5, t)

    # ================================================================
    # SLIDE 6: TECHNICAL ARCHITECTURE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip(s, t)

    add_tag(s, Inches(0.5), Inches(0.25), "Technical Architecture", t["blue"])
    add_text(s, Inches(0.5), Inches(0.55), Inches(12), Inches(0.5),
             "Built for enterprise scale, security, and observability",
             font_size=30, color=t["heading"], bold=True)

    # ── LAYER 1: Frontend ──
    layer1 = add_rounded_rect(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.7), t)
    layer1.line.color.rgb = t["blue"]
    add_text(s, Inches(0.6), Inches(1.25), Inches(1.8), Inches(0.6),
             "Frontend", font_size=13, color=t["blue"], bold=True)
    fe_items = ["React / TypeScript", "Vite", "30+ Pages", "Real-time Dashboards", "Semantic Search UI"]
    for i, item in enumerate(fe_items):
        x = Inches(2.5) + Inches(i * 2.1)
        add_text(s, x, Inches(1.3), Inches(2.0), Inches(0.5),
                 item, font_size=10, color=t["heading"], bold=False, alignment=PP_ALIGN.CENTER)

    # ── LAYER 2: API ──
    layer2 = add_rounded_rect(s, Inches(0.4), Inches(2.1), Inches(12.5), Inches(0.7), t)
    layer2.line.color.rgb = t["purple_lt"]
    add_text(s, Inches(0.6), Inches(2.15), Inches(1.8), Inches(0.6),
             "API Layer", font_size=13, color=t["purple_lt"], bold=True)
    api_items = ["FastAPI (async)", "315+ Endpoints", "41 Routers", "Multi-Tenant Isolation", "JWT + RBAC"]
    for i, item in enumerate(api_items):
        x = Inches(2.5) + Inches(i * 2.1)
        add_text(s, x, Inches(2.2), Inches(2.0), Inches(0.5),
                 item, font_size=10, color=t["heading"], bold=False, alignment=PP_ALIGN.CENTER)

    # ── LAYER 3: Services + AI Agents (side by side) ──
    # Services block (left)
    svc = add_rounded_rect(s, Inches(0.4), Inches(3.0), Inches(6.0), Inches(1.9), t)
    svc.line.color.rgb = t["emerald"]
    add_text(s, Inches(0.6), Inches(3.05), Inches(2.5), Inches(0.3),
             "Services (36 modules)", font_size=13, color=t["emerald"], bold=True)

    svc_items = [
        ("Contract Pipeline", "Upload, parse, chunk,\nembed, index"),
        ("Governance Bridge", "7 automations: org, relationship,\nKPI, health, improvements"),
        ("Golden Dataset", "Few-shot learning, verified\nextractions, quality metrics"),
        ("Auto-Link Detector", "6-signal weighted matching,\ncontract family detection"),
    ]
    for i, (title, desc) in enumerate(svc_items):
        col = i % 2; row = i // 2
        x = Inches(0.6) + Inches(col * 2.95)
        y = Inches(3.4) + Inches(row * 0.72)
        add_dot(s, x, y + Inches(0.04), t["emerald"], Inches(0.08))
        add_text(s, x + Inches(0.15), y, Inches(1.3), Inches(0.25),
                 title, font_size=10, color=t["heading"], bold=True)
        add_text(s, x + Inches(0.15), y + Inches(0.22), Inches(2.6), Inches(0.45),
                 desc, font_size=8, color=t["body"])

    # AI Agents block (right)
    ai = add_rounded_rect(s, Inches(6.7), Inches(3.0), Inches(6.2), Inches(1.9), t)
    ai.line.color.rgb = t["orange"]
    add_text(s, Inches(6.9), Inches(3.05), Inches(3.5), Inches(0.3),
             "9 AI Agents (GPT-4o + Agent Squad)", font_size=13, color=t["orange"], bold=True)

    agent_items = [
        "Metadata Extraction", "Risk Detection", "Clause Extraction",
        "Obligation Tracking", "SLA Extraction", "Renewal Monitoring",
        "Schema Extraction", "Intent Router", "Contract Q&A (RAG)",
    ]
    for i, name in enumerate(agent_items):
        col = i % 3; row = i // 3
        x = Inches(6.9) + Inches(col * 2.0)
        y = Inches(3.45) + Inches(row * 0.45)
        add_dot(s, x, y + Inches(0.04), t["orange"], Inches(0.08))
        add_text(s, x + Inches(0.15), y, Inches(1.8), Inches(0.25),
                 name, font_size=9, color=t["heading"])

    # ── LAYER 4: Data stores ──
    # PostgreSQL
    pg = add_rounded_rect(s, Inches(0.4), Inches(5.1), Inches(4.0), Inches(0.85), t)
    pg.line.color.rgb = t["blue"]
    add_text(s, Inches(0.6), Inches(5.15), Inches(3.6), Inches(0.3),
             "PostgreSQL", font_size=13, color=t["blue"], bold=True)
    add_text(s, Inches(0.6), Inches(5.45), Inches(3.6), Inches(0.4),
             "55+ tables  \u2022  Row-level tenant isolation  \u2022  Alembic migrations",
             font_size=9, color=t["body"])

    # ChromaDB
    ch = add_rounded_rect(s, Inches(4.65), Inches(5.1), Inches(3.7), Inches(0.85), t)
    ch.line.color.rgb = t["purple_lt"]
    add_text(s, Inches(4.85), Inches(5.15), Inches(3.3), Inches(0.3),
             "ChromaDB", font_size=13, color=t["purple_lt"], bold=True)
    add_text(s, Inches(4.85), Inches(5.45), Inches(3.3), Inches(0.4),
             "Vector embeddings  \u2022  Semantic search  \u2022  RAG pipeline",
             font_size=9, color=t["body"])

    # Integrations
    integ = add_rounded_rect(s, Inches(8.6), Inches(5.1), Inches(4.3), Inches(0.85), t)
    integ.line.color.rgb = t["emerald"]
    add_text(s, Inches(8.8), Inches(5.15), Inches(3.9), Inches(0.3),
             "Integrations", font_size=13, color=t["emerald"], bold=True)
    add_text(s, Inches(8.8), Inches(5.45), Inches(3.9), Inches(0.4),
             "ServiceNow  \u2022  Salesforce  \u2022  Teams  \u2022  REST API  \u2022  Webhooks",
             font_size=9, color=t["body"])

    # ── LAYER 5: Cross-cutting (bottom bar) ──
    # Security
    sec = add_rounded_rect(s, Inches(0.4), Inches(6.15), Inches(6.0), Inches(0.6), t)
    sec.line.color.rgb = t["red"]
    add_text(s, Inches(0.6), Inches(6.2), Inches(1.2), Inches(0.45),
             "Security", font_size=11, color=t["red"], bold=True)
    add_text(s, Inches(1.7), Inches(6.22), Inches(4.5), Inches(0.45),
             "SOC 2-ready  \u2022  GDPR-aware  \u2022  AES-256  \u2022  RBAC  \u2022  Multi-tenant  \u2022  Audit Trails",
             font_size=9, color=t["body"])

    # Observability
    obs = add_rounded_rect(s, Inches(6.7), Inches(6.15), Inches(6.2), Inches(0.6), t)
    obs.line.color.rgb = t["amber"]
    add_text(s, Inches(6.9), Inches(6.2), Inches(1.8), Inches(0.45),
             "Observability", font_size=11, color=t["amber"], bold=True)
    add_text(s, Inches(8.3), Inches(6.22), Inches(4.4), Inches(0.45),
             "Langfuse  \u2022  Full LLM tracing  \u2022  Cost per contract  \u2022  OpenTelemetry",
             font_size=9, color=t["body"])

    add_slide_number(s, 6, t)

    # ================================================================
    # SLIDE 7: RISK DETECTION
    # ================================================================
    build_deep_dive(
        slide_num=7,
        tag_text="Deep Dive", tag_color=t["red"],
        title="Risk Detection: see what others miss",
        subtitle=(
            "Every contract carries risk \u2014 hidden indemnities, uncapped liabilities, "
            "ambiguous termination clauses. Evaluetor\u2019s AI scores and categorises "
            "risk automatically so your team can focus on what matters most."
        ),
        left_title="How Risk Analysis Works",
        left_subtitle="Fully automated from upload to alert:",
        pipeline_steps=[
            ("Document Parsing", "Text extracted from PDF, DOCX, or scanned images via OCR", t["purple_lt"]),
            ("Clause Classification", "AI identifies 30+ clause types across the full document", t["blue"]),
            ("Risk Categorisation", "Each clause scored across 10 risk categories", t["orange"]),
            ("Severity Scoring", "4-level scale: Low, Medium, High, Critical with confidence", t["red"]),
            ("Portfolio Alerts", "High-risk items surfaced on dashboards with recommended actions", t["pink"]),
        ],
        right_title="Key Capabilities",
        feature_cards=[
            ("10 Risk Categories",
             "Indemnification, liability caps, IP assignment, termination, "
             "confidentiality, non-compete, data protection, insurance, "
             "warranty, and force majeure \u2014 each scored with weighted "
             "severity from 0\u2013100.",
             t["red"]),
            ("Cross-Contract Anomaly Detection",
             "Knowledge graph compares terms across your portfolio and flags "
             "deviations from your norms. Example: \u201cThis vendor usually has "
             "30-day payment terms \u2014 this contract has 45.\u201d",
             t["orange"]),
            ("Portfolio Risk Dashboard",
             "Aggregate risk posture across all contracts with drill-down by "
             "vendor, contract type, business unit, or risk category. Trend "
             "analysis shows whether your portfolio risk is improving over time.",
             t["purple_lt"]),
        ],
    )

    # ================================================================
    # SLIDE 8: OBLIGATION TRACKING
    # ================================================================
    build_deep_dive(
        slide_num=8,
        tag_text="Deep Dive", tag_color=t["emerald"],
        title="Obligation Tracking: never miss a commitment",
        subtitle=(
            "Missed obligations cost money, damage relationships, and create compliance exposure. "
            "Evaluetor extracts every obligation from every contract and tracks them "
            "through their full lifecycle \u2014 with owners, deadlines, and evidence."
        ),
        left_title="The Obligation Lifecycle",
        left_subtitle="From extraction to compliance proof:",
        pipeline_steps=[
            ("AI Extraction", "Obligations identified with deadlines, parties, and consequences", t["purple_lt"]),
            ("Classification", "7 types: delivery, reporting, payment, regulatory, performance, notice, compliance", t["emerald"]),
            ("Owner Assignment", "Each obligation assigned to a responsible person with escalation path", t["orange"]),
            ("Status Tracking", "RAG indicators (Red / Amber / Green) updated as deadlines approach", t["amber"]),
            ("Evidence & Audit", "Compliance evidence attached for audit trails and reporting", t["pink"]),
        ],
        right_title="Key Capabilities",
        feature_cards=[
            ("RAG Status Dashboard",
             "At-a-glance compliance visibility across your entire portfolio. "
             "Red flags overdue items, amber shows approaching deadlines, "
             "green confirms on-track \u2014 filterable by vendor, business unit, "
             "or contract type.",
             t["emerald"]),
            ("Configurable Deadline Alerts",
             "Automated notifications at configurable lead times before each "
             "deadline. Alerts route to obligation owners via email, Slack, "
             "or Microsoft Teams \u2014 with escalation when deadlines are missed.",
             t["orange"]),
            ("Compliance Reporting",
             "Generate compliance reports by date range, vendor, or obligation "
             "type. Track fulfilment rates over time with trend analysis. "
             "Attach evidence documents for audit readiness.",
             t["blue"]),
        ],
        use_orange_strip=True,
    )

    # ================================================================
    # SLIDE 9: SLA MONITORING
    # ================================================================
    build_deep_dive(
        slide_num=9,
        tag_text="Deep Dive", tag_color=t["blue"],
        title="SLA Monitoring: recover what you\u2019re owed",
        subtitle=(
            "SLA breaches happen constantly \u2014 the challenge is detecting them. "
            "Evaluetor extracts SLA terms, compares them against actual performance, "
            "and calculates every service credit you\u2019re entitled to."
        ),
        left_title="From Contract to Credit Recovery",
        left_subtitle="Automated end-to-end:",
        pipeline_steps=[
            ("SLA Extraction", "AI identifies metrics, targets, measurement periods, and penalties", t["purple_lt"]),
            ("Performance Ingestion", "Actual performance data imported via API or manual entry", t["blue"]),
            ("Threshold Comparison", "Contracted targets compared against actuals in real time", t["emerald"]),
            ("Breach Detection", "Violations classified by severity: Minor, Moderate, Major, Critical", t["orange"]),
            ("Credit Calculation", "Service credits and penalties computed automatically per contract terms", t["amber"]),
        ],
        right_title="Key Capabilities",
        feature_cards=[
            ("9 SLA Metric Types",
             "Uptime, Response Time, Resolution Time, Throughput, Error Rate, "
             "Availability, Latency, MTTR, and Custom metrics \u2014 each with "
             "targets, thresholds, and measurement windows extracted by AI.",
             t["blue"]),
            ("Service Credit Recovery",
             "Automatic calculation of credits owed based on breach severity "
             "and contract terms. Track recovery status from detection through "
             "claim to resolution. Customers report 3\u20135x more credits found.",
             t["emerald"]),
            ("Vendor Scorecards",
             "A\u2013F grade vendor ratings combining SLA compliance rate, breach "
             "frequency, financial impact, and trend direction. Compare vendor "
             "performance across your portfolio with historical benchmarking.",
             t["orange"]),
        ],
    )

    # ================================================================
    # SLIDE 10: RELATIONSHIP GOVERNANCE
    # ================================================================
    build_deep_dive(
        slide_num=10,
        tag_text="Key Differentiator", tag_color=t["pink"],
        title="Relationship Governance: only in Evaluetor",
        subtitle=(
            "Most CLM platforms stop at the contract. Evaluetor goes further \u2014 automatically "
            "connecting contract data to the business relationships they support. "
            "No other CLM vendor offers this capability."
        ),
        left_title="The Automated Governance Bridge",
        left_subtitle="Contract upload triggers a fully automated pipeline:",
        pipeline_steps=[
            ("Contract Upload", "AI extracts parties, terms, SLAs, and obligations", t["orange"]),
            ("Organizations", "Counterparties auto-linked to your org directory", t["blue"]),
            ("Relationships", "Business relationships created with health scoring", t["emerald"]),
            ("KPIs & Scoring", "Performance metrics tracked with dual-perspective scoring", t["amber"]),
            ("Improvement Points", "Gaps automatically surface with action recommendations", t["pink"]),
        ],
        right_title="What Makes This Unique",
        feature_cards=[
            ("Perception Gap Analysis",
             "Dual-perspective scoring reveals misalignment between how you see "
             "vendor performance and how they see it \u2014 with severity classification "
             "that drives targeted improvement.",
             t["pink"]),
            ("Composite Health Scoring",
             "Relationship health calculated from Risk (30%), SLA Compliance (40%), "
             "and Obligation Fulfillment (30%) \u2014 a single score that tells you "
             "where each relationship truly stands.",
             t["emerald"]),
            ("Automated Improvement Tracking",
             "When perception gaps exceed thresholds, improvement points are "
             "auto-generated with ownership, priority, and action items. "
             "From insight to action without manual intervention.",
             t["orange"]),
        ],
        use_orange_strip=True,
    )

    # ================================================================
    # SLIDE 11: WHY EVALUETOR (COMPARISON)
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip_orange(s, t)

    add_tag(s, Inches(0.8), Inches(0.6), "Why Evaluetor", t["purple_lt"])
    add_text(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "AI-native. Not AI-bolted-on.", font_size=40,
             color=t["heading"], bold=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.5),
             "Legacy CLM tools added AI as an afterthought. We built intelligence into the foundation.",
             font_size=16, color=t["body"])

    headers = ["Capability", "Evaluetor", "DocuSign CLM", "Icertis", "Ironclad"]
    rows_data = [
        ["Time to Value",            "Hours",           "6\u201312 months", "12\u201318 months", "3\u20136 months"],
        ["AI Extraction",            "9 native agents",  "Limited",         "Add-on",            "Limited"],
        ["SLA Monitoring",           "Built-in",         "No",              "Add-on",            "No"],
        ["Obligation Tracking",      "AI-extracted",     "Manual",          "Manual",            "Limited"],
        ["Relationship Governance",  "Built-in",         "No",              "No",                "No"],
        ["Knowledge Graph",          "Built-in",         "No",              "No",                "No"],
        ["Unlimited Users",          "\u2713 Yes",       "\u2717",          "\u2717",            "\u2717"],
        ["LLM Observability",        "Langfuse",         "No",              "No",                "No"],
    ]

    col_widths = [Inches(2.8), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]
    table = s.shapes.add_table(len(rows_data) + 1, 5,
                                Inches(0.8), Inches(2.6),
                                Inches(12.0), Inches(4.8))
    tbl = table.table

    for ci in range(5):
        tbl.columns[ci].width = col_widths[ci]

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = t["orange"] if ci == 1 else t["muted"]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = t["tbl_hdr"]

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            p.font.bold = ci == 1
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            if ci == 1:
                p.font.color.rgb = t["orange"]
            elif ci == 0:
                p.font.color.rgb = t["label"]
            else:
                p.font.color.rgb = t["muted"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = t["tbl_alt"] if ri % 2 == 0 else t["bg"]

    add_slide_number(s, 11, t)

    # ================================================================
    # SLIDE 12: RESULTS & NEXT STEPS
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, t)
    accent_strip(s, t)

    add_tag(s, Inches(0.8), Inches(0.6), "Results & Next Steps", t["emerald"])
    add_text(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
             "Measurable impact from day one", font_size=40,
             color=t["heading"], bold=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.5),
             "Immediate, quantifiable results \u2014 not months of implementation before you see value.",
             font_size=16, color=t["body"])

    metrics = [
        ("90%", "Less time\nreviewing contracts", t["purple_lt"]),
        ("3\u20135x", "More SLA credits\nrecovered", t["orange"]),
        ("95%", "Fewer missed\ndeadlines", t["red"]),
        ("$2M+", "Credits recovered\nacross customers", t["emerald"]),
    ]

    for i, (val, label, accent) in enumerate(metrics):
        x = Inches(0.8) + Inches(i * 3.1)
        add_rounded_rect(s, x, Inches(2.7), Inches(2.7), Inches(1.8), t)
        add_text(s, x + Inches(0.3), Inches(2.9), Inches(2.1), Inches(0.8),
                 val, font_size=42, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), Inches(3.65), Inches(2.1), Inches(0.7),
                 label, font_size=13, color=t["body"], alignment=PP_ALIGN.CENTER)

    teams = [
        ("For Legal", "Days of contract review compressed to minutes. "
         "Automated risk scoring across 10 categories. Portfolio-wide semantic search.",
         t["purple_lt"]),
        ("For Procurement", "Real-time SLA monitoring. Detect breaches "
         "the moment they happen. Recover every service credit owed.",
         t["orange"]),
        ("For Finance", "Track every obligation and deadline. Detect "
         "auto-renewal traps before they trigger. Understand true financial exposure.",
         t["emerald"]),
    ]

    for i, (title, desc, accent) in enumerate(teams):
        x = Inches(0.8) + Inches(i * 4.0)
        add_rounded_rect(s, x, Inches(5.0), Inches(3.6), Inches(1.6), t)
        add_dot(s, x + Inches(0.3), Inches(5.3), accent)
        add_text(s, x + Inches(0.55), Inches(5.2), Inches(2.8), Inches(0.35),
                 title, font_size=16, color=t["heading"], bold=True)
        add_text(s, x + Inches(0.3), Inches(5.65), Inches(3.0), Inches(0.85),
                 desc, font_size=12, color=t["body"])

    cta = add_rounded_rect(s, Inches(4.65), Inches(6.85), Inches(4.0), Inches(0.55), t)
    cta.fill.solid()
    cta.fill.fore_color.rgb = t["orange"]
    cta.line.fill.background()
    add_text(s, Inches(4.65), Inches(6.87), Inches(4.0), Inches(0.55),
             "Let\u2019s schedule a live demo  \u2192", font_size=18,
             color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
             alignment=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(7.0), Inches(4), Inches(0.4),
             "evaluetor.com", font_size=10, color=t["muted"])
    add_slide_number(s, 12, t)

    # ── Save ──
    prs.save(output_path)
    print(f"Saved: {output_path}")


# ───────────────────────────────────────────────────────────────
# GENERATE BOTH DECKS
# ───────────────────────────────────────────────────────────────

build_deck(DARK_THEME,  "/Users/jjayaraj/workspaces/studios/clm/sales/Evaluetor-Sales-Pitch-Deck.pptx")
build_deck(LIGHT_THEME, "/Users/jjayaraj/workspaces/studios/clm/sales/Evaluetor-Sales-Pitch-Deck-Light.pptx")
