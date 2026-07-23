"""Internal sales training playbook — 10-slide PPTX.

Inverts the prospect-facing deck: this one is FOR the seller, not the
prospect. Captures Steven Visser's slide-8 CFO negotiation playbook plus
ideal deal structures from slide 7.

Output: sales/Evaluetor-Sales-Training-Playbook.pptx
Run:    python3 generate_sales_training.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Palette (matches prospect-slider, Evaluetor brand) ───────────
PAPER      = RGBColor(0xF7, 0xF6, 0xF3)
PAPER_2    = RGBColor(0xEF, 0xED, 0xE8)
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x0E, 0x0E, 0x10)
INK_2      = RGBColor(0x2D, 0x2F, 0x36)
INK_3      = RGBColor(0x5B, 0x5F, 0x66)
INK_4      = RGBColor(0xA0, 0xA4, 0xAB)
RULE       = RGBColor(0x1A, 0x1A, 0x1C)
RULE_2     = RGBColor(0xDD, 0xDB, 0xD5)
ORANGE     = RGBColor(0xE9, 0x4E, 0x1B)
ORANGE_2   = RGBColor(0xC7, 0x3E, 0x0F)
RED        = RGBColor(0xB9, 0x1C, 0x1C)
GREEN      = RGBColor(0x15, 0x80, 0x3D)

FONT_SANS = "Calibri"
FONT_MONO = "Consolas"
FONT_ITAL = "Cambria"

SW = Inches(13.333); SH = Inches(7.5)
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]

# ─── Helpers ──────────────────────────────────────────────────────
def set_bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_width: sh.line.width = line_width
    sh.shadow.inherit = False
    return sh

def hline(slide, x, y, w, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(1))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    return sh

def tx(slide, left, top, width, height, text,
       font=FONT_SANS, size=14, color=INK, bold=False, italic=False,
       align=PP_ALIGN.LEFT, spacing=0, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing

    if isinstance(text, list):
        for i, item in enumerate(text):
            if isinstance(item, tuple):
                t, opts = item
            else:
                t, opts = item, {}
            r = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
            r.text = t
            f = r.font
            f.name = opts.get('font', font)
            f.size = Pt(opts.get('size', size))
            f.bold = opts.get('bold', bold)
            f.italic = opts.get('italic', italic)
            f.color.rgb = opts.get('color', color)
            if spacing:
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(spacing * 100)))
    else:
        p.text = text
        for r in p.runs:
            f = r.font
            f.name = font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
            if spacing:
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(spacing * 100)))
    return box

def slide_chrome(slide, head_label, foot_label, bg=PAPER):
    set_bg(slide, bg)
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.28), Inches(0.28), ORANGE_2)
    tx(slide, Inches(0.55), Inches(0.43), Inches(0.28), Inches(0.28),
       "E", font=FONT_MONO, size=11, color=SURFACE, bold=True, align=PP_ALIGN.CENTER)
    tx(slide, Inches(0.92), Inches(0.42), Inches(3.5), Inches(0.3),
       "Evaluetor", font=FONT_SANS, size=14, color=INK, bold=True)
    tx(slide, Inches(4.5), Inches(0.42), Inches(3), Inches(0.3),
       "·  SALES TRAINING", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
    tx(slide, Inches(8.0), Inches(0.43), Inches(4.8), Inches(0.3),
       head_label, font=FONT_MONO, size=9, color=INK_3, spacing=1.6,
       align=PP_ALIGN.RIGHT)
    hline(slide, Inches(0.55), Inches(0.85), Inches(12.23), RULE)
    hline(slide, Inches(0.55), Inches(6.95), Inches(12.23), RULE_2)
    tx(slide, Inches(0.55), Inches(7.05), Inches(7), Inches(0.25),
       "INTERNAL · NOT FOR PROSPECTS",
       font=FONT_MONO, size=9, color=ORANGE_2, bold=True, spacing=1.6)
    tx(slide, Inches(7.0), Inches(7.05), Inches(5.78), Inches(0.25),
       foot_label,
       font=FONT_MONO, size=9, color=INK_3, spacing=1.4, align=PP_ALIGN.RIGHT)

def section_label(slide, x, y, text):
    tx(slide, x, y, Inches(8), Inches(0.3),
       text, font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)

def objection_card(slide, x, y, w, h, num, objection, response_runs):
    """Render one objection card."""
    rect(slide, x, y, w, h, SURFACE, line_color=RULE_2)
    rect(slide, x, y, Inches(0.18), h, ORANGE_2)
    tx(slide, x + Inches(0.32), y + Inches(0.18), Inches(2), Inches(0.3),
       f"OBJECTION {num}", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
    tx(slide, x + Inches(0.32), y + Inches(0.48), w - Inches(0.5), Inches(0.5),
       objection, font=FONT_SANS, size=17, color=INK, bold=True, italic=True)
    tx(slide, x + Inches(0.32), y + Inches(1.05), w - Inches(0.5), h - Inches(1.3),
       response_runs, line_spacing=1.4)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 · COVER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "COVER", "Slide 01 — Cover", bg=PAPER)

rect(s, Inches(0.7), Inches(1.6), Inches(0.85), Inches(0.85), ORANGE_2)
tx(s, Inches(0.7), Inches(1.7), Inches(0.85), Inches(0.7),
   "E", font=FONT_MONO, size=32, color=SURFACE, bold=True, align=PP_ALIGN.CENTER)

tx(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.2),
   "Sales Training", font=FONT_SANS, size=68, color=INK, bold=False)
tx(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
   "How to run an Evaluetor conversation.",
   font=FONT_ITAL, size=26, color=ORANGE_2, italic=True)

tx(s, Inches(0.7), Inches(4.7), Inches(11), Inches(0.5),
   "A negotiation playbook for selling AI-native execution intelligence to procurement-led enterprises.",
   font=FONT_SANS, size=15, color=INK_2)

hline(s, Inches(0.7), Inches(5.7), Inches(10), RULE_2)
tx(s, Inches(0.7), Inches(5.85), Inches(3), Inches(0.25),
   "AUDIENCE", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(0.7), Inches(6.15), Inches(3), Inches(0.3),
   "Internal sales team", font=FONT_SANS, size=13, color=INK, bold=True)
tx(s, Inches(4.2), Inches(5.85), Inches(3), Inches(0.25),
   "BASED ON", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(4.2), Inches(6.15), Inches(4), Inches(0.3),
   "Steven Visser, May 2026", font=FONT_SANS, size=13, color=INK, bold=True)
tx(s, Inches(8.5), Inches(5.85), Inches(3), Inches(0.25),
   "USE WITH", font=FONT_MONO, size=8.5, color=INK_3, spacing=1.6)
tx(s, Inches(8.5), Inches(6.15), Inches(4), Inches(0.3),
   "Excel value calculator", font=FONT_SANS, size=13, color=INK, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 · CFO MINDSET + OPENING POSITION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "02 / MINDSET + OPENING", "Slide 02 — CFO mindset + opening position", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Understand the CFO. ", {"size": 36, "color": INK, "font": FONT_SANS}),
    ("Then open right.", {"size": 36, "color": INK_2, "font": FONT_ITAL, "italic": True})])

# Left column: CFO mindset
section_label(s, Inches(0.55), Inches(2.05), "1.  WHAT THE CFO ACTUALLY CARES ABOUT")
hline(s, Inches(0.55), Inches(2.4), Inches(5.8), RULE_2)

tx(s, Inches(0.55), Inches(2.55), Inches(5.8), Inches(0.35),
   "Not these:", font=FONT_SANS, size=13, color=INK_3, italic=True)
for i, item in enumerate(["AI agents", "CLM features", "Dashboards"]):
    rect(s, Inches(0.55), Inches(2.95) + i*Inches(0.36), Inches(5.8), Inches(0.32), SURFACE, line_color=RULE_2)
    tx(s, Inches(0.7), Inches(3.02) + i*Inches(0.36), Inches(5.6), Inches(0.3),
       f"❌  {item}", font=FONT_SANS, size=13, color=INK_3)

tx(s, Inches(0.55), Inches(4.25), Inches(5.8), Inches(0.35),
   "These:", font=FONT_SANS, size=13, color=ORANGE_2, italic=True, bold=True)
for i, item in enumerate(["Cash impact", "Forecast reliability", "Risk avoidance", "Audit defensibility"]):
    rect(s, Inches(0.55), Inches(4.65) + i*Inches(0.36), Inches(5.8), Inches(0.32), PAPER_2, line_color=ORANGE_2)
    tx(s, Inches(0.7), Inches(4.72) + i*Inches(0.36), Inches(5.6), Inches(0.3),
       f"✓  {item}", font=FONT_SANS, size=13, color=INK, bold=True)

# Right column: Opening line
section_label(s, Inches(6.7), Inches(2.05), "2.  YOUR FIRST STATEMENT — NEVER SKIP IT")
hline(s, Inches(6.7), Inches(2.4), Inches(6.1), RULE_2)

rect(s, Inches(6.7), Inches(2.6), Inches(6.1), Inches(2.4), INK)
rect(s, Inches(6.7), Inches(2.6), Inches(0.18), Inches(2.4), ORANGE)
tx(s, Inches(6.95), Inches(2.85), Inches(5.6), Inches(0.4),
   "OPEN WITH", font=FONT_MONO, size=10, color=ORANGE, bold=True, spacing=1.6)
tx(s, Inches(6.95), Inches(3.2), Inches(5.7), Inches(1.7),
   [("We typically identify ", {"size": 17, "color": PAPER, "font": FONT_SANS}),
    ("€1–2M of recoverable value", {"size": 17, "color": PAPER, "font": FONT_SANS, "bold": True}),
    (" per €50M of contract volume. Our fee is ", {"size": 17, "color": PAPER, "font": FONT_SANS}),
    ("a fraction of that.", {"size": 17, "color": ORANGE, "font": FONT_ITAL, "italic": True})],
   line_spacing=1.35)

tx(s, Inches(6.7), Inches(5.2), Inches(6.1), Inches(0.35),
   "WHY IT WORKS", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
hline(s, Inches(6.7), Inches(5.55), Inches(6.1), RULE_2)
tx(s, Inches(6.7), Inches(5.7), Inches(6.1), Inches(1.2),
   [("It instantly reframes the conversation: ", {"size": 13, "color": INK_2}),
    ("price → investment, cost → return multiple", {"size": 13, "color": INK, "bold": True}),
    (".  The CFO is now thinking about ROI before they think about budget category.", {"size": 13, "color": INK_2})],
   line_spacing=1.45)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 · 3-LAYER PRICE DEFENSE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "03 / PRICE DEFENSE", "Slide 03 — Three-layer price defense", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Three layers of price defense. ", {"size": 36, "color": INK, "font": FONT_SANS}),
    ("In this order.", {"size": 36, "color": INK_2, "font": FONT_ITAL, "italic": True})])

tx(s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.4),
   "Never anchor on the price tag. Anchor on three things, in sequence.",
   font=FONT_SANS, size=14, color=INK_3)

layers = [
    ("LAYER 01", "ROI multiple", "This is a 10–25× ROI decision.",
     "CFO mental model:\n< 3× → reject  ·  5–10× → consider  ·  > 10× → approve",
     "Always lead here."),
    ("LAYER 02", "Cost of inaction", "Doing nothing costs you €1–2M per year.",
     "Make concrete: missed SLA credits  ·  surprise auto-renewals  ·  pricing leakage  ·  scope-creep unbilled.",
     "Use against status quo bias."),
    ("LAYER 03", "Budget reframing", "This should sit next to procurement savings or audit recovery — not IT tooling.",
     "❌ 'CLM software'\n✓ 'Value recovery / cost optimisation'",
     "Re-categorise the spend.")
]

card_w = Inches(4.1); card_h = Inches(4.4)
for i, (label, title, line, body, footer) in enumerate(layers):
    x = Inches(0.55) + i * (card_w + Inches(0.1))
    rect(s, x, Inches(2.45), card_w, card_h, SURFACE, line_color=RULE_2)
    rect(s, x, Inches(2.45), card_w, Inches(0.06), ORANGE_2)
    tx(s, x + Inches(0.3), Inches(2.65), card_w - Inches(0.5), Inches(0.3),
       label, font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
    tx(s, x + Inches(0.3), Inches(3.0), card_w - Inches(0.5), Inches(0.4),
       title, font=FONT_SANS, size=20, color=INK, bold=True)
    tx(s, x + Inches(0.3), Inches(3.5), card_w - Inches(0.5), Inches(1.0),
       line, font=FONT_ITAL, size=15, color=INK_2, italic=True, line_spacing=1.3)
    hline(s, x + Inches(0.3), Inches(4.85), card_w - Inches(0.6), RULE_2)
    tx(s, x + Inches(0.3), Inches(5.0), card_w - Inches(0.5), Inches(1.4),
       body, font=FONT_SANS, size=12, color=INK_2, line_spacing=1.4)
    tx(s, x + Inches(0.3), Inches(6.45), card_w - Inches(0.5), Inches(0.3),
       footer, font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.4)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 · OBJECTION 1: "This is expensive"
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "04 / OBJECTION 1", "Slide 04 — Objection 1: Expensive", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Objection 01.  ", {"size": 36, "color": ORANGE_2, "font": FONT_SANS, "bold": True}),
    ('"This is expensive."', {"size": 36, "color": INK, "font": FONT_ITAL, "italic": True})])

section_label(s, Inches(0.55), Inches(2.1), "YOUR RESPONSE")
hline(s, Inches(0.55), Inches(2.45), Inches(12.23), RULE_2)
tx(s, Inches(0.55), Inches(2.6), Inches(12), Inches(0.6),
   '"Relative to what?"',
   font=FONT_ITAL, size=32, color=INK, italic=True, bold=True)

tx(s, Inches(0.55), Inches(3.3), Inches(12), Inches(0.4),
   "Then offer three benchmarks they cannot ignore:",
   font=FONT_SANS, size=14, color=INK_3)

benchmarks = [
    ("€80k", "vs €1.5M recoverable annual value"),
    ("€80k", "vs ONE missed SLA claim event"),
    ("€80k", "vs one incorrectly handled renewal"),
]
for i, (lhs, rhs) in enumerate(benchmarks):
    y = Inches(3.85) + i * Inches(0.55)
    rect(s, Inches(0.55), y, Inches(12.23), Inches(0.45), SURFACE, line_color=RULE_2)
    tx(s, Inches(0.8), y + Inches(0.1), Inches(1.5), Inches(0.3),
       lhs, font=FONT_SANS, size=18, color=ORANGE_2, bold=True)
    tx(s, Inches(2.4), y + Inches(0.13), Inches(0.5), Inches(0.3),
       "vs", font=FONT_MONO, size=11, color=INK_3, spacing=1.4)
    tx(s, Inches(3.0), y + Inches(0.13), Inches(9), Inches(0.3),
       rhs, font=FONT_SANS, size=14, color=INK)

section_label(s, Inches(0.55), Inches(5.75), "CLOSING LINE")
hline(s, Inches(0.55), Inches(6.1), Inches(12.23), RULE_2)
rect(s, Inches(0.55), Inches(6.25), Inches(12.23), Inches(0.55), INK)
tx(s, Inches(0.85), Inches(6.36), Inches(11.7), Inches(0.4),
   '"You\'re comparing it to software. You should compare it to lost cash."',
   font=FONT_ITAL, size=17, color=PAPER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 · OBJECTION 2: "We already have a CLM"
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "05 / OBJECTION 2", "Slide 05 — Objection 2: Already have a CLM", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Objection 02.  ", {"size": 36, "color": ORANGE_2, "font": FONT_SANS, "bold": True}),
    ('"We already have a CLM."', {"size": 36, "color": INK, "font": FONT_ITAL, "italic": True})])

section_label(s, Inches(0.55), Inches(2.1), "YOUR RESPONSE")
hline(s, Inches(0.55), Inches(2.45), Inches(12.23), RULE_2)

tx(s, Inches(0.55), Inches(2.6), Inches(12), Inches(0.6),
   [('"You already have a contract repository. ', {"size": 24, "color": INK, "font": FONT_SANS}),
    ('What you don\'t have is execution visibility."', {"size": 24, "color": ORANGE_2, "font": FONT_ITAL, "italic": True, "bold": True})])

# Two-column: incumbent vs us
y0 = Inches(3.7)
col_w = Inches(6.0)

rect(s, Inches(0.55), y0, col_w, Inches(1.9), PAPER_2, line_color=RULE_2)
tx(s, Inches(0.75), y0 + Inches(0.2), col_w - Inches(0.4), Inches(0.3),
   "INCUMBENT CLM", font=FONT_MONO, size=10, color=INK_3, bold=True, spacing=1.6)
tx(s, Inches(0.75), y0 + Inches(0.55), col_w - Inches(0.4), Inches(0.4),
   "Storage + workflow", font=FONT_SANS, size=20, color=INK, bold=True)
tx(s, Inches(0.75), y0 + Inches(1.0), col_w - Inches(0.4), Inches(1.0),
   "Drafting → routing → eSignature → repository.  Designed in the document era.",
   font=FONT_SANS, size=13, color=INK_2, line_spacing=1.4)

rect(s, Inches(6.78), y0, col_w, Inches(1.9), INK)
rect(s, Inches(6.78), y0, Inches(0.18), Inches(1.9), ORANGE)
tx(s, Inches(7.05), y0 + Inches(0.2), col_w - Inches(0.4), Inches(0.3),
   "EVALUETOR", font=FONT_MONO, size=10, color=ORANGE, bold=True, spacing=1.6)
tx(s, Inches(7.05), y0 + Inches(0.55), col_w - Inches(0.4), Inches(0.4),
   "€ recovery + risk detection", font=FONT_SANS, size=20, color=PAPER, bold=True)
tx(s, Inches(7.05), y0 + Inches(1.0), col_w - Inches(0.4), Inches(1.0),
   "Reads what they signed.  Watches it perform.  Surfaces leakage.  Built model-first.",
   font=FONT_SANS, size=13, color=PAPER, line_spacing=1.4)

# Kill-switch line
rect(s, Inches(0.55), Inches(6.0), Inches(12.23), Inches(0.7), INK)
rect(s, Inches(0.55), Inches(6.0), Inches(0.18), Inches(0.7), ORANGE)
tx(s, Inches(0.85), Inches(6.05), Inches(11.7), Inches(0.3),
   "KILL SWITCH", font=FONT_MONO, size=9, color=ORANGE, bold=True, spacing=1.6)
tx(s, Inches(0.85), Inches(6.32), Inches(11.7), Inches(0.4),
   '"If your current CLM can show € recovered / € lost monthly for every contract, you don\'t need us."',
   font=FONT_ITAL, size=15, color=PAPER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 · OBJECTION 3: "Integration will be complex"
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "06 / OBJECTION 3", "Slide 06 — Objection 3: Integration", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Objection 03.  ", {"size": 36, "color": ORANGE_2, "font": FONT_SANS, "bold": True}),
    ('"Integration will be complex."', {"size": 36, "color": INK, "font": FONT_ITAL, "italic": True})])

section_label(s, Inches(0.55), Inches(2.1), "YOUR RESPONSE  ·  TWO MOVES")
hline(s, Inches(0.55), Inches(2.45), Inches(12.23), RULE_2)

# Step 1
tx(s, Inches(0.55), Inches(2.7), Inches(2), Inches(0.3),
   "MOVE 01", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
tx(s, Inches(2.0), Inches(2.7), Inches(11), Inches(0.3),
   "Acknowledge — don't fight.", font=FONT_SANS, size=15, color=INK_3, italic=True)

rect(s, Inches(0.55), Inches(3.1), Inches(12.23), Inches(0.7), SURFACE, line_color=ORANGE_2)
rect(s, Inches(0.55), Inches(3.1), Inches(0.18), Inches(0.7), ORANGE_2)
tx(s, Inches(0.85), Inches(3.25), Inches(11.7), Inches(0.5),
   '"Correct — this is where most of the value sits."',
   font=FONT_ITAL, size=17, color=INK, italic=True, bold=True)

# Step 2
tx(s, Inches(0.55), Inches(4.15), Inches(2), Inches(0.3),
   "MOVE 02", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
tx(s, Inches(2.0), Inches(4.15), Inches(11), Inches(0.3),
   "Turn it on its head.", font=FONT_SANS, size=15, color=INK_3, italic=True)

rect(s, Inches(0.55), Inches(4.55), Inches(12.23), Inches(0.95), INK)
rect(s, Inches(0.55), Inches(4.55), Inches(0.18), Inches(0.95), ORANGE)
tx(s, Inches(0.85), Inches(4.7), Inches(11.7), Inches(0.7),
   '"That\'s exactly why the ROI exists. If the data were already connected, the value wouldn\'t be leaking."',
   font=FONT_ITAL, size=17, color=PAPER, italic=True, line_spacing=1.3)

# How to act on this
section_label(s, Inches(0.55), Inches(5.85), "AND THEN")
hline(s, Inches(0.55), Inches(6.2), Inches(12.23), RULE_2)
tx(s, Inches(0.55), Inches(6.35), Inches(12), Inches(0.6),
   "Limit initial scope → ServiceNow OR one ERP, not both.  Expand after proof.  ROI within 90 days on Tier 02.",
   font=FONT_SANS, size=14, color=INK_2, line_spacing=1.4)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 · OBJECTION 4: "We need proof"
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "07 / OBJECTION 4", "Slide 07 — Objection 4: Proof", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Objection 04.  ", {"size": 36, "color": ORANGE_2, "font": FONT_SANS, "bold": True}),
    ('"We need proof before committing."', {"size": 36, "color": INK, "font": FONT_ITAL, "italic": True})])

section_label(s, Inches(0.55), Inches(2.1), "YOUR RESPONSE  ·  REMOVE THE RISK")
hline(s, Inches(0.55), Inches(2.45), Inches(12.23), RULE_2)

tx(s, Inches(0.55), Inches(2.6), Inches(12), Inches(0.5),
   '"Agreed. That\'s why we start with 5 contracts."',
   font=FONT_ITAL, size=24, color=INK, italic=True, bold=True)

# What you find in 5 contracts
tx(s, Inches(0.55), Inches(3.6), Inches(12), Inches(0.3),
   "What you typically surface from 5 contracts in 30 minutes:",
   font=FONT_SANS, size=14, color=INK_3)

findings = [
    ("Missed obligations", "with concrete deadlines and owners"),
    ("Renewal risks", "auto-renewal traps, missed notice periods"),
    ("SLA gaps", "uncalculated credits, breach exposure"),
    ("Indemnity / liability risks", "uncapped, broad, or one-sided terms"),
]
for i, (lhs, rhs) in enumerate(findings):
    y = Inches(4.05) + i * Inches(0.42)
    rect(s, Inches(0.55), y, Inches(12.23), Inches(0.35), SURFACE, line_color=RULE_2)
    rect(s, Inches(0.55), y, Inches(0.06), Inches(0.35), ORANGE_2)
    tx(s, Inches(0.85), y + Inches(0.07), Inches(4), Inches(0.3),
       lhs, font=FONT_SANS, size=13, color=INK, bold=True)
    tx(s, Inches(5.0), y + Inches(0.07), Inches(8), Inches(0.3),
       rhs, font=FONT_SANS, size=13, color=INK_3)

# Closing
section_label(s, Inches(0.55), Inches(5.85), "CLOSING LINE  ·  REMOVES ALL BUYING RISK")
hline(s, Inches(0.55), Inches(6.2), Inches(12.23), RULE_2)
rect(s, Inches(0.55), Inches(6.35), Inches(12.23), Inches(0.5), INK)
tx(s, Inches(0.85), Inches(6.42), Inches(11.7), Inches(0.4),
   '"If we don\'t find value, you don\'t proceed."',
   font=FONT_ITAL, size=17, color=PAPER, italic=True, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 · OBJECTION 5: "Can you discount?"
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "08 / OBJECTION 5", "Slide 08 — Objection 5: Discount", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Objection 05.  ", {"size": 36, "color": ORANGE_2, "font": FONT_SANS, "bold": True}),
    ('"Can you discount?"', {"size": 36, "color": INK, "font": FONT_ITAL, "italic": True})])

section_label(s, Inches(0.55), Inches(2.05), "NEVER DISCOUNT FIRST.  TRADE.")
hline(s, Inches(0.55), Inches(2.4), Inches(12.23), RULE_2)

# Step 1 anchor
tx(s, Inches(0.55), Inches(2.6), Inches(2), Inches(0.3),
   "STEP 01", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
tx(s, Inches(2.0), Inches(2.6), Inches(11), Inches(0.3),
   "Anchor value first.", font=FONT_SANS, size=13, color=INK_3, italic=True)
tx(s, Inches(0.55), Inches(2.95), Inches(12), Inches(0.4),
   '"If we reduce scope, we reduce value recovery."',
   font=FONT_ITAL, size=18, color=INK, italic=True)

# Step 2 trade-off table
tx(s, Inches(0.55), Inches(3.75), Inches(2), Inches(0.3),
   "STEP 02", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
tx(s, Inches(2.0), Inches(3.75), Inches(11), Inches(0.3),
   "Offer a trade.  You give one thing, get one thing.",
   font=FONT_SANS, size=13, color=INK_3, italic=True)

trades = [
    ("You give", "You get"),
    ("Lower price", "Longer contract (2–3 yrs)"),
    ("Lower price", "Named reference case"),
    ("Lower price", "Faster decision / no procurement loop"),
]
row_h = Inches(0.42)
for i, (l, r) in enumerate(trades):
    y = Inches(4.15) + i * row_h
    is_header = i == 0
    bg = INK if is_header else SURFACE
    border = INK if is_header else RULE_2
    rect(s, Inches(0.55), y, Inches(6.1), row_h, bg, line_color=border)
    rect(s, Inches(6.7), y, Inches(6.1), row_h, bg, line_color=border)
    text_color = PAPER if is_header else INK_2
    font_size = 11 if is_header else 14
    font_weight = True if is_header else False
    tx(s, Inches(0.75), y + Inches(0.1), Inches(5.7), Inches(0.3),
       l, font=FONT_MONO if is_header else FONT_SANS, size=font_size,
       color=text_color, bold=font_weight,
       spacing=1.6 if is_header else 0)
    tx(s, Inches(6.9), y + Inches(0.1), Inches(5.7), Inches(0.3),
       r, font=FONT_MONO if is_header else FONT_SANS, size=font_size,
       color=text_color, bold=font_weight,
       spacing=1.6 if is_header else 0)

# Step 3
tx(s, Inches(0.55), Inches(6.1), Inches(2), Inches(0.3),
   "STEP 03", font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
tx(s, Inches(2.0), Inches(6.1), Inches(11), Inches(0.3),
   "Controlled concession.", font=FONT_SANS, size=13, color=INK_3, italic=True)
tx(s, Inches(0.55), Inches(6.42), Inches(12), Inches(0.4),
   '"We can optimise pricing if we structure this as a multi-year engagement."',
   font=FONT_ITAL, size=15, color=INK, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 · NEGOTIATION LEVERS + CLOSING STRATEGY
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "09 / LEVERS + CLOSING", "Slide 09 — Negotiation levers + closing", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Levers you can pull. ", {"size": 36, "color": INK, "font": FONT_SANS}),
    ("And how to close.", {"size": 36, "color": INK_2, "font": FONT_ITAL, "italic": True})])

# Left: levers
section_label(s, Inches(0.55), Inches(2.05), "WHAT YOU CAN MOVE  ·  NEVER €/VALUE")
hline(s, Inches(0.55), Inches(2.4), Inches(5.8), RULE_2)

levers = [
    ("Scope",        "# contracts  ·  business units  ·  modules"),
    ("Time",         "pilot vs full rollout  ·  staged expansion"),
    ("Integrations", "delay connectors → lower price"),
    ("Commitment",   "1 yr → higher price  ·  3 yr → lower price"),
]
for i, (k, v) in enumerate(levers):
    y = Inches(2.6) + i * Inches(0.6)
    rect(s, Inches(0.55), y, Inches(5.8), Inches(0.55), SURFACE, line_color=RULE_2)
    rect(s, Inches(0.55), y, Inches(0.18), Inches(0.55), ORANGE_2)
    tx(s, Inches(0.85), y + Inches(0.06), Inches(1.5), Inches(0.3),
       k, font=FONT_SANS, size=13, color=INK, bold=True)
    tx(s, Inches(0.85), y + Inches(0.3), Inches(5.4), Inches(0.3),
       v, font=FONT_SANS, size=12, color=INK_3)

# Right: closing strategy
section_label(s, Inches(6.7), Inches(2.05), "CLOSING STRATEGY  ·  4 STEPS")
hline(s, Inches(6.7), Inches(2.4), Inches(6.1), RULE_2)

steps = [
    ("01", "Summarise value",  '"We\'ve identified €X risk / missed value."'),
    ("02", "Restate ROI",      '"You invest €Y to recover €X."'),
    ("03", "Reduce risk",      '"We start with a focused scope."'),
    ("04", "Call to action",   '"Do you want visibility into this value now or next year?"'),
]
for i, (n, k, v) in enumerate(steps):
    y = Inches(2.6) + i * Inches(0.95)
    tx(s, Inches(6.7), y, Inches(0.4), Inches(0.3),
       n, font=FONT_MONO, size=12, color=ORANGE_2, bold=True)
    tx(s, Inches(7.1), y, Inches(5), Inches(0.3),
       k, font=FONT_SANS, size=14, color=INK, bold=True)
    tx(s, Inches(7.1), y + Inches(0.32), Inches(5.5), Inches(0.5),
       v, font=FONT_ITAL, size=13, color=INK_2, italic=True, line_spacing=1.3)

tx(s, Inches(6.7), Inches(6.5), Inches(6), Inches(0.4),
   "Make delay the real cost.",
   font=FONT_MONO, size=11, color=ORANGE_2, bold=True, spacing=1.6)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 · IDEAL DEAL STRUCTURES (TARGETS)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
slide_chrome(s, "10 / DEAL STRUCTURES", "Slide 10 — Ideal deal structures", bg=PAPER)

tx(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.7),
   [("Aim for these. ", {"size": 36, "color": INK, "font": FONT_SANS}),
    ("Negotiate around them.", {"size": 36, "color": INK_2, "font": FONT_ITAL, "italic": True})])

tx(s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.4),
   "Every deal must support a >10× ROI narrative.  Below that, restructure the scope before discounting.",
   font=FONT_SANS, size=14, color=INK_3)

# Two deal structure cards
deals = [
    ("MID-ENTERPRISE  ·  TARGET",
     "€20M – €100M ACVUM",
     [("Base platform fee",  "€80k"),
      ("Connectors (1–2)",   "€30k"),
      ("Implementation",     "€40k")],
     "€150k",
     "Total Year 1",
     "€110k",
     "Recurring"),
    ("ENTERPRISE  ·  TARGET",
     "€100M – €500M ACVUM",
     [("Base platform fee",  "€150k – €250k"),
      ("Connectors (2–3)",   "€50k"),
      ("Implementation",     "€80k")],
     "€280k – €380k",
     "Total Year 1",
     "€200k – €300k",
     "Recurring")
]

card_w = Inches(6.0); card_h = Inches(4.4)
for i, (label, scope, items, y1, y1l, rec, recl) in enumerate(deals):
    x = Inches(0.55) + i * (card_w + Inches(0.2))
    rect(s, x, Inches(2.45), card_w, card_h, SURFACE, line_color=RULE_2)
    rect(s, x, Inches(2.45), card_w, Inches(0.06), ORANGE_2)
    tx(s, x + Inches(0.3), Inches(2.65), card_w - Inches(0.5), Inches(0.3),
       label, font=FONT_MONO, size=10, color=ORANGE_2, bold=True, spacing=1.6)
    tx(s, x + Inches(0.3), Inches(3.0), card_w - Inches(0.5), Inches(0.4),
       scope, font=FONT_SANS, size=20, color=INK, bold=True)

    hline(s, x + Inches(0.3), Inches(3.55), card_w - Inches(0.6), RULE_2)
    for j, (k, v) in enumerate(items):
        ry = Inches(3.7) + j * Inches(0.38)
        tx(s, x + Inches(0.3), ry, card_w - Inches(2.0), Inches(0.3),
           k, font=FONT_SANS, size=13, color=INK_2)
        tx(s, x + card_w - Inches(2.0), ry, Inches(1.7), Inches(0.3),
           v, font=FONT_MONO, size=12, color=INK, bold=True, align=PP_ALIGN.RIGHT)

    hline(s, x + Inches(0.3), Inches(5.0), card_w - Inches(0.6), RULE_2)

    # Year 1 total
    tx(s, x + Inches(0.3), Inches(5.2), card_w - Inches(0.5), Inches(0.3),
       y1l, font=FONT_MONO, size=10, color=INK_3, bold=True, spacing=1.4)
    tx(s, x + Inches(0.3), Inches(5.5), card_w - Inches(0.5), Inches(0.5),
       y1, font=FONT_SANS, size=24, color=ORANGE_2, bold=True)

    # Recurring
    tx(s, x + Inches(0.3), Inches(6.15), card_w - Inches(0.5), Inches(0.3),
       recl, font=FONT_MONO, size=10, color=INK_3, bold=True, spacing=1.4)
    tx(s, x + Inches(0.3), Inches(6.45), card_w - Inches(0.5), Inches(0.4),
       rec, font=FONT_SANS, size=18, color=INK, bold=True)


# ─── Save ────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Evaluetor-Sales-Training-Playbook.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
