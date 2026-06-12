"""Generate Evaluetor Commercial Presentation — Light theme, 10 slides, tight layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Palette ───
BG       = RGBColor(0xFF, 0xFF, 0xFF)
CARD     = RGBColor(0xF9, 0xFA, 0xFB)
CARD_ALT = RGBColor(0xF3, 0xF4, 0xF6)
HEADING  = RGBColor(0x11, 0x18, 0x27)
BODY     = RGBColor(0x4B, 0x55, 0x63)
MUTED    = RGBColor(0x6B, 0x72, 0x80)
BORDER   = RGBColor(0xE5, 0xE7, 0xEB)
CONN     = RGBColor(0xD1, 0xD5, 0xDB)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
ORANGE   = RGBColor(0xEA, 0x76, 0x0D)
PINK     = RGBColor(0xDB, 0x27, 0x77)
RED      = RGBColor(0xDC, 0x26, 0x26)
EMERALD  = RGBColor(0x05, 0x96, 0x69)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
CYAN     = RGBColor(0x08, 0x91, 0xB2)
TEAL     = RGBColor(0x0D, 0x94, 0x88)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL = 14
SW = Inches(13.333); SH = Inches(7.5)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]

# ─── Helpers ───
def set_bg(s): f = s.background.fill; f.solid(); f.fore_color.rgb = BG
def bar(s):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.06))
    sh.line.fill.background(); sh.fill.solid(); sh.fill.fore_color.rgb = BLUE
def rect(s, l, t, w, h, fill=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.color.rgb = BORDER; sh.line.width = Pt(1)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; return sh
def tx(s, l, t, w, h, text, sz=18, clr=HEADING, bold=False, align=PP_ALIGN.LEFT):
    b = s.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.font.name = "Calibri"
    p.alignment = align; return b
def sn(s, n):
    tx(s, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
       f"{n}/{TOTAL}", sz=10, clr=MUTED, align=PP_ALIGN.RIGHT)
def dot(s, x, y, clr, sz=Inches(0.12)):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    d.line.fill.background(); d.fill.solid(); d.fill.fore_color.rgb = clr
def conn_line(s, x, y, h):
    c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.04), h)
    c.line.fill.background(); c.fill.solid(); c.fill.fore_color.rgb = CONN
def footer(s):
    tx(s, Inches(0.5), Inches(7.08), Inches(1.5), Inches(0.3),
       "13-4-2026", sz=8, clr=MUTED, bold=True)
    tx(s, Inches(4.0), Inches(7.08), Inches(4), Inches(0.3),
       "AI-Powered Contract Intelligence", sz=8, clr=MUTED, bold=True,
       align=PP_ALIGN.CENTER)
    tx(s, Inches(10.5), Inches(7.05), Inches(2.7), Inches(0.35),
       "Evaluetor\u00a9 AI", sz=13, clr=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    # Separator line
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.015))
    ln.line.fill.background(); ln.fill.solid(); ln.fill.fore_color.rgb = BORDER

def deep_dive(num, tag_clr, title, subtitle,
              left_title, left_sub, steps, right_title, cards):
    s = prs.slides.add_slide(BL); set_bg(s); bar(s)
    # Title
    tx(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
       title, sz=28, clr=HEADING, bold=True)
    # Subtitle block — compact
    rect(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(1.3), CARD)
    tx(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.1),
       subtitle, sz=14, clr=BODY, align=PP_ALIGN.CENTER)
    # Left pipeline
    tx(s, Inches(0.5), Inches(2.35), Inches(5.5), Inches(0.3),
       left_title, sz=15, clr=HEADING, bold=True)
    tx(s, Inches(0.5), Inches(2.65), Inches(5.5), Inches(0.25),
       left_sub, sz=10, clr=MUTED)
    for i, (st, sd, ac) in enumerate(steps):
        y = Inches(3.0) + Inches(i * 0.82)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y,
                                Inches(0.32), Inches(0.32))
        c.line.fill.background(); c.fill.solid(); c.fill.fore_color.rgb = ac
        tx(s, Inches(0.7), y + Inches(0.02), Inches(0.32), Inches(0.3),
           str(i+1), sz=11, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx(s, Inches(1.2), y + Inches(0.02), Inches(2.0), Inches(0.3),
           st, sz=12, clr=HEADING, bold=True)
        tx(s, Inches(3.2), y + Inches(0.04), Inches(3.3), Inches(0.3),
           sd, sz=10, clr=BODY)
        if i < len(steps) - 1:
            conn_line(s, Inches(0.84), y + Inches(0.35), Inches(0.47))
    # Right cards
    rx = Inches(6.8)
    tx(s, rx, Inches(2.35), Inches(6), Inches(0.3),
       right_title, sz=15, clr=HEADING, bold=True)
    for i, (ft, fd, ac) in enumerate(cards):
        y = Inches(2.8) + Inches(i * 1.4)
        rect(s, rx, y, Inches(6.0), Inches(1.25), CARD)
        dot(s, rx + Inches(0.2), y + Inches(0.18), ac)
        tx(s, rx + Inches(0.45), y + Inches(0.1),
           Inches(5.3), Inches(0.28), ft, sz=13, clr=HEADING, bold=True)
        tx(s, rx + Inches(0.2), y + Inches(0.4),
           Inches(5.6), Inches(0.8), fd, sz=10, clr=BODY)
    footer(s); sn(s, num)


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
   "Reducing Contract Value Leakage Through Active Contract Management",
   sz=30, clr=HEADING, bold=True, align=PP_ALIGN.CENTER)

# Full-width card
rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6), CARD)

# Left side — key message
tx(s, Inches(1.0), Inches(1.5), Inches(5.5), Inches(1.5),
   "Evaluetor\u00a9 AI\ndelivering contract\nbenefits effectively",
   sz=32, clr=HEADING, bold=True)

tx(s, Inches(1.0), Inches(3.3), Inches(5.5), Inches(1.2),
   "AI-native contract intelligence that makes your contracts "
   "transparent, measurable, and actionable \u2014 in minutes, not months.",
   sz=15, clr=BODY)

# Highlight stat
rect(s, Inches(1.0), Inches(4.7), Inches(5.2), Inches(1.5), CARD_ALT)
tx(s, Inches(1.3), Inches(4.85), Inches(4.6), Inches(0.5),
   "Organizations lose 8\u201311% of contract value post-signing",
   sz=16, clr=ORANGE, bold=True)
tx(s, Inches(1.3), Inches(5.35), Inches(4.6), Inches(0.7),
   "due to poor execution and monitoring. Active management "
   "reduces losses by 30\u201340%.",
   sz=13, clr=BODY)

# Right side — three callout cards stacked
callouts = [
    ("Contracts are negotiated carefully, then stored away. "
     "The intent behind the terms is lost.", ORANGE),
    ("Enterprises struggle to surface what\u2019s actually inside "
     "their contract portfolio.", RED),
    ("Risks, obligations, SLAs, and renewal terms often go "
     "completely unmonitored.", BLUE),
]
for i, (txt, accent) in enumerate(callouts):
    y = Inches(1.4) + Inches(i * 1.65)
    rect(s, Inches(7.0), y, Inches(5.5), Inches(1.4), CARD_ALT)
    dot(s, Inches(7.3), y + Inches(0.2), accent, Inches(0.14))
    tx(s, Inches(7.6), y + Inches(0.15), Inches(4.7), Inches(1.1),
       txt, sz=14, clr=HEADING, bold=True)

footer(s); sn(s, 1)


# ================================================================
# SLIDE 2: VALUE LEAKAGES (BUY & SELL SIDE)
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
   "Where Contract Value Leaks — Buy Side and Sell Side",
   sz=26, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.35),
   "70–80% of contract value leakage happens post-signature — not at the negotiation table.",
   sz=13, clr=ORANGE, bold=True)

# Left column: Top 10 leakages
rect(s, Inches(0.5), Inches(1.3), Inches(6.3), Inches(5.4), CARD)
tx(s, Inches(0.75), Inches(1.4), Inches(5.9), Inches(0.35),
   "Top 10 Contract Value Leakages", sz=15, clr=HEADING, bold=True)
leakages = [
    "Non-compliance with contract terms (off-contract buying / under-billing)",
    "Poor pricing & discount control",
    "Missed rebates, incentives & credits",
    "Weak contract lifecycle management (visibility, access)",
    "Scope creep & unmanaged changes",
    "SLA underperformance not enforced or penalised",
    "Ineffective renewals & auto-renewals",
    "Volume commitment & forecast inaccuracies",
    "Legal–commercial misalignment",
    "No clear contract ownership & accountability",
]
for i, item in enumerate(leakages):
    y = Inches(1.85) + Inches(i * 0.46)
    # Numbered circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y + Inches(0.05),
                            Inches(0.3), Inches(0.3))
    c.line.fill.background(); c.fill.solid(); c.fill.fore_color.rgb = RED
    tx(s, Inches(0.75), y + Inches(0.07), Inches(0.3), Inches(0.28),
       str(i+1), sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tx(s, Inches(1.15), y + Inches(0.08), Inches(5.55), Inches(0.4),
       item, sz=11, clr=BODY)

# Right column: Operating Model Principles
rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.4), CARD_ALT)
tx(s, Inches(7.25), Inches(1.4), Inches(5.4), Inches(0.35),
   "Operating Model Principles (Target State)", sz=15, clr=HEADING, bold=True)
principles = [
    ("Contract ≠ Document",
     "Contract = data + obligations + risk", PURPLE),
    ("Post-signature is the value zone",
     "70–80% of leakage happens after signing", ORANGE),
    ("Business owns value; Legal owns integrity",
     "Clear separation of accountability", BLUE),
    ("One contract truth across systems",
     "Single source of contract data", EMERALD),
    ("No unmanaged deviations",
     "Every change tracked and approved", AMBER),
    ("Value realisation is measured, not assumed",
     "Track actual outcomes against negotiated terms", TEAL),
]
for i, (title, desc, accent) in enumerate(principles):
    y = Inches(1.85) + Inches(i * 0.78)
    dot(s, Inches(7.25), y + Inches(0.12), accent, Inches(0.14))
    tx(s, Inches(7.55), y + Inches(0.02), Inches(5.1), Inches(0.3),
       title, sz=12, clr=HEADING, bold=True)
    tx(s, Inches(7.55), y + Inches(0.32), Inches(5.1), Inches(0.42),
       desc, sz=10, clr=BODY)

footer(s); sn(s, 2)


# ================================================================
# SLIDE 3: WHY ACTIVE CONTRACT MANAGEMENT
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
   "Why Active Contract Management Is a Commercial Imperative",
   sz=28, clr=HEADING, bold=True)

points = [
    ("Value Leakage Post-Signing",
     "Organizations lose 8\u201311% of contract value after signing due to poor execution "
     "and monitoring. The problem isn\u2019t negotiation \u2014 it\u2019s what happens after the ink dries.",
     ORANGE),
    ("Systemic Commercial Failure",
     "Value loss comes from lack of ownership and enforcement, not negotiation. Without "
     "active tracking, negotiated terms never translate into realized financial outcomes.",
     RED),
    ("Active Management Benefits",
     "Active contract management assigns ownership, tracks obligations, and captures value "
     "\u2014 reducing losses by 30\u201340%. Every contract becomes a managed financial instrument.",
     EMERALD),
    ("Profit Enablement Capability",
     "Active contract management turns negotiated agreements into realized financial results "
     "and predictable cash flow. Contracts stop being static documents and start driving revenue.",
     BLUE),
]
for i, (title, desc, accent) in enumerate(points):
    y = Inches(1.0) + Inches(i * 1.55)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(1.35), CARD)
    dot(s, Inches(0.85), y + Inches(0.2), accent, Inches(0.16))
    tx(s, Inches(1.2), y + Inches(0.12), Inches(5.0), Inches(0.35),
       title, sz=16, clr=HEADING, bold=True)
    tx(s, Inches(1.2), y + Inches(0.5), Inches(11.3), Inches(0.75),
       desc, sz=13, clr=BODY)

footer(s); sn(s, 3)


# ================================================================
# SLIDE 4: CONTRACT LIFECYCLE
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
   "Supporting the Contract Lifecycle", sz=28, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.4),
   "Most CLM tools focus on pre-signature. Evaluetor focuses on what happens after \u2014 "
   "where the real value is won or lost.",
   sz=14, clr=BODY)

# Three phases — taller cards
phases = [
    ("Contracting", "Negotiation, drafting,\nand execution", MUTED, False),
    ("Execution", "Obligation tracking, SLA monitoring,\nrisk management, relationship governance", BLUE, True),
    ("Close / Renewal", "Renewal decisions,\nrenegotiation, exit", MUTED, False),
]
for i, (title, desc, clr, hl) in enumerate(phases):
    x = Inches(0.5) + Inches(i * 4.2)
    bg_c = RGBColor(0xEF, 0xF6, 0xFF) if hl else CARD
    c = rect(s, x, Inches(1.4), Inches(3.9), Inches(1.8), bg_c)
    if hl: c.line.color.rgb = BLUE; c.line.width = Pt(2)
    tx(s, x + Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.35),
       title, sz=18, clr=BLUE if hl else MUTED, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.2), Inches(1.95), Inches(3.5), Inches(0.8),
       desc, sz=12, clr=HEADING if hl else BODY, align=PP_ALIGN.CENTER)
    if hl:
        tx(s, x + Inches(0.2), Inches(2.75), Inches(3.5), Inches(0.3),
           "Evaluetor\u00a9 AI", sz=14, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)
# Arrows
for i in range(2):
    x = Inches(4.4) + Inches(i * 4.2)
    tx(s, x, Inches(2.0), Inches(0.4), Inches(0.4),
       "\u2192", sz=24, clr=CONN, bold=True, align=PP_ALIGN.CENTER)

# Challenges — 2 rows x 4
tx(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.3),
   "Challenges that arise during execution:", sz=14, clr=HEADING, bold=True)
challenges = [
    "Scope changes", "Dissatisfaction", "Complaints", "Disputes",
    "Legislation changes", "New requirements", "Add-on requests", "SLA breaches",
]
for i, ch in enumerate(challenges):
    col = i % 4; row = i // 4
    x = Inches(0.5) + Inches(col * 3.15)
    y = Inches(3.9) + Inches(row * 0.6)
    rect(s, x, y, Inches(2.95), Inches(0.45), CARD)
    dot(s, x + Inches(0.12), y + Inches(0.14), ORANGE, Inches(0.1))
    tx(s, x + Inches(0.32), y + Inches(0.08), Inches(2.5), Inches(0.3),
       ch, sz=12, clr=HEADING)

# Bottom CTA
rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3), RGBColor(0xEF, 0xF6, 0xFF))
tx(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.9),
   "Evaluetor\u2019s AI agents monitor all of this automatically \u2014 detecting risks, "
   "tracking obligations, and alerting stakeholders before value is lost.",
   sz=16, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)

footer(s); sn(s, 4)


# ================================================================
# SLIDE 5: CONNECTORS — EXTERNAL DATA SOURCES
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
   "Connected to where conditions are measurable",
   sz=26, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.35),
   "Evaluetor connects contracts with the outside world — comparing contracted terms "
   "against real-time data from the systems where work actually happens.",
   sz=13, clr=BODY)

# Three connector cards
connectors = [
    ("Volumes", "ERP",
     "SAP, Oracle, NetSuite",
     "Compare actual purchase volumes against committed volumes. "
     "Detect under-buying and trigger rebate claims automatically.", PURPLE),
    ("SLA", "ServiceNow",
     "ITSM / ticketing systems",
     "Ingest incident response times, resolution metrics, and uptime. "
     "Detect SLA breaches the moment they happen, calculate credits.", BLUE),
    ("FX", "ECB Website",
     "Central bank / market data",
     "Pull live FX rates for currency-indexed contracts. "
     "Detect when contractual FX thresholds are crossed.", EMERALD),
]
for i, (metric, system, source, desc, accent) in enumerate(connectors):
    x = Inches(0.5) + Inches(i * 4.28)
    rect(s, x, Inches(1.45), Inches(4.1), Inches(3.3), CARD)
    # Header band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x, Inches(1.45), Inches(4.1), Inches(0.85))
    band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = accent
    tx(s, x + Inches(0.25), Inches(1.55), Inches(3.6), Inches(0.4),
       f"{metric} @ {system}", sz=18, clr=WHITE, bold=True)
    tx(s, x + Inches(0.25), Inches(1.95), Inches(3.6), Inches(0.3),
       source, sz=11, clr=WHITE)
    # Body
    tx(s, x + Inches(0.25), Inches(2.5), Inches(3.6), Inches(2.0),
       desc, sz=12, clr=BODY)

# Bottom pipeline — Upload → AI → Dashboards → Action
rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6), CARD_ALT)
tx(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.3),
   "Continuous Control Cycle by Agentic AI Agents",
   sz=13, clr=HEADING, bold=True, align=PP_ALIGN.CENTER)
pipe = [
    ("01", "Upload",     "Contracts: MSA, schedules, SOWs", PURPLE),
    ("02", "AI Extracts", "Structured & analysed",          ORANGE),
    ("03", "Dashboards",  "With external data comparison",  TEAL),
    ("04", "Take Action", "Agentic AI continuous control",  BLUE),
]
for i, (num, title, desc, accent) in enumerate(pipe):
    x = Inches(0.7) + Inches(i * 3.0)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(5.55),
                            Inches(0.5), Inches(0.5))
    c.line.color.rgb = accent; c.line.width = Pt(2)
    c.fill.solid(); c.fill.fore_color.rgb = BG
    tx(s, x, Inches(5.6), Inches(0.5), Inches(0.4),
       num, sz=12, clr=accent, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.6), Inches(5.55), Inches(2.3), Inches(0.3),
       title, sz=13, clr=HEADING, bold=True)
    tx(s, x + Inches(0.6), Inches(5.85), Inches(2.3), Inches(0.5),
       desc, sz=10, clr=BODY)

footer(s); sn(s, 5)


# ================================================================
# SLIDE 6: ONE PLATFORM
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.15), Inches(3.2), Inches(0.4),
   "Evaluetor\u00a9 AI", sz=18, clr=BLUE, bold=True)
tx(s, Inches(3.5), Inches(0.15), Inches(8), Inches(0.4),
   "One platform. Complete intelligence.", sz=18, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.25),
   "From upload to intelligence in minutes", sz=12, clr=ORANGE, bold=True,
   align=PP_ALIGN.CENTER)

# Pipeline row
rect(s, Inches(0.4), Inches(0.9), Inches(12.5), Inches(2.5), CARD)
step_colors = [PURPLE, ORANGE, TEAL, BLUE]
step_data = [
    ("01", "Upload", "Drop any PDF, DOCX, or batch\nupload a ZIP. OCR handles scans."),
    ("02", "AI Extracts", "9 agents extract metadata,\nclauses, obligations, SLAs, risks."),
    ("03", "Dashboards", "Real-time dashboards surface\nrisks, deadlines, and vendor scores."),
    ("04", "Take Action", "Claim credits, meet deadlines,\nrenegotiate, govern relationships."),
]
for i, (num, title, desc) in enumerate(step_data):
    x = Inches(0.6) + Inches(i * 3.1)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(1.1),
                            Inches(0.65), Inches(0.65))
    c.line.color.rgb = step_colors[i]; c.line.width = Pt(2)
    c.fill.solid(); c.fill.fore_color.rgb = BG
    tx(s, x + Inches(0.55), Inches(1.17), Inches(0.65), Inches(0.5),
       num, sz=16, clr=step_colors[i], bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), Inches(1.85), Inches(1.6), Inches(0.3),
       title, sz=13, clr=HEADING, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x, Inches(2.15), Inches(1.8), Inches(0.7),
       desc, sz=10, clr=BODY, align=PP_ALIGN.CENTER)

# Tagline
tx(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.25),
   "No templates. No pre-configuration. No training required. Just upload and go.",
   sz=11, clr=MUTED, align=PP_ALIGN.CENTER)

# 6 capability cards — 2x3 grid, tighter
cap_data = [
    ("Contract Intelligence", "AI extracts metadata, 30+ clause types, and risks across "
     "10 categories. Full semantic search across entire portfolio.", PURPLE),
    ("Obligation Tracking", "RAG status indicators, owner assignment, escalation paths, "
     "and automated deadline alerts. Never miss a commitment.", EMERALD),
    ("SLA Monitoring", "Real-time breach detection, service credit calculation, and vendor "
     "scorecards. Know when performance falls short.", BLUE),
    ("Renewal Management", "Auto-renewal detection, notice period tracking, and AI-powered "
     "renewal recommendations. Stay ahead of every date.", ORANGE),
    ("Relationship Governance", "KPI perception gap analysis, vendor health scores, satisfaction "
     "surveys, and improvement tracking. Only available in Evaluetor.", PINK),
    ("Enterprise Security", "SOC 2-ready architecture, GDPR-aware, AES-256 encryption, "
     "role-based access, multi-tenant isolation, full audit trails.", CYAN),
]
for i, (title, desc, accent) in enumerate(cap_data):
    col = i % 3; row = i // 3
    x = Inches(0.4) + Inches(col * 4.25)
    y = Inches(3.85) + Inches(row * 1.55)
    rect(s, x, y, Inches(4.05), Inches(1.4), CARD)
    dot(s, x + Inches(0.15), y + Inches(0.15), accent, Inches(0.1))
    tx(s, x + Inches(0.35), y + Inches(0.08), Inches(3.5), Inches(0.28),
       title, sz=12, clr=HEADING, bold=True)
    tx(s, x + Inches(0.15), y + Inches(0.4), Inches(3.75), Inches(0.9),
       desc, sz=10, clr=BODY)

footer(s); sn(s, 6)


# ================================================================
# SLIDES 7-10: DEEP DIVES
# ================================================================
deep_dive(7, EMERALD,
    "Obligation Tracking: never miss a commitment",
    "Missed obligations cost money, damage relationships, and create compliance exposure. "
    "Evaluetor extracts every obligation from every contract and tracks them through their "
    "full lifecycle \u2014 with owners, deadlines, and evidence.",
    "The Obligation Lifecycle", "From extraction to compliance proof:",
    [("AI Extraction", "Obligations identified with deadlines, parties, and consequences", PURPLE),
     ("Classification", "7 types: delivery, reporting, payment, regulatory, performance, notice, compliance", EMERALD),
     ("Owner Assignment", "Each obligation assigned to a responsible person with escalation path", ORANGE),
     ("Status Tracking", "RAG indicators (Red / Amber / Green) updated as deadlines approach", AMBER),
     ("Evidence & Audit", "Compliance evidence attached for audit trails and reporting", BLUE)],
    "Key Capabilities",
    [("RAG Status Dashboard",
      "At-a-glance compliance visibility across your entire portfolio. Red flags overdue items, "
      "amber shows approaching deadlines, green confirms on-track \u2014 filterable by vendor, "
      "business unit, or contract type.", EMERALD),
     ("Configurable Deadline Alerts",
      "Automated notifications at configurable lead times before each deadline. Alerts route to "
      "obligation owners via email, Slack, or Microsoft Teams \u2014 with escalation when deadlines "
      "are missed.", ORANGE),
     ("Compliance Reporting",
      "Generate compliance reports by date range, vendor, or obligation type. Track fulfilment "
      "rates over time with trend analysis. Attach evidence documents for audit readiness.", BLUE)])

deep_dive(8, RED,
    "Risk Detection: see what others miss",
    "Every contract carries risk \u2014 hidden indemnities, uncapped liabilities, ambiguous termination "
    "clauses. Evaluetor\u2019s AI scores and categorises risk automatically so your team can focus on "
    "what matters most.",
    "How Risk Analysis Works", "Fully automated from upload to alert:",
    [("Document Parsing", "Text extracted from PDF, DOCX, or scanned images via OCR", PURPLE),
     ("Clause Classification", "AI identifies 30+ clause types across the full document", BLUE),
     ("Risk Categorisation", "Each clause scored across 10 risk categories", ORANGE),
     ("Severity Scoring", "4-level scale: Low, Medium, High, Critical with confidence", RED),
     ("Portfolio Alerts", "High-risk items surfaced on dashboards with recommended actions", PINK)],
    "Key Capabilities",
    [("10 Risk Categories",
      "Indemnification, liability caps, IP assignment, termination, confidentiality, non-compete, "
      "data protection, insurance, warranty, and force majeure \u2014 each scored with weighted severity "
      "from 0\u2013100.", RED),
     ("Cross-Contract Anomaly Detection",
      "Knowledge graph compares terms across your portfolio and flags deviations from your norms. "
      "Example: \u201cThis vendor usually has 30-day payment terms \u2014 this contract has 45.\u201d", ORANGE),
     ("Portfolio Risk Dashboard",
      "Aggregate risk posture across all contracts with drill-down by vendor, contract type, "
      "business unit, or risk category. Trend analysis shows whether portfolio risk is improving.", PURPLE)])

deep_dive(9, BLUE,
    "SLA Monitoring: recover what you\u2019re owed",
    "SLA breaches happen constantly \u2014 the challenge is detecting them. Evaluetor extracts SLA terms, "
    "compares them against actual performance, and calculates every service credit you\u2019re entitled to.",
    "From Contract to Credit Recovery", "Automated end-to-end:",
    [("SLA Extraction", "AI identifies metrics, targets, measurement periods, and penalties", PURPLE),
     ("Performance Ingestion", "Actual performance data imported via API or manual entry", BLUE),
     ("Threshold Comparison", "Contracted targets compared against actuals in real time", EMERALD),
     ("Breach Detection", "Violations classified by severity: Minor, Moderate, Major, Critical", ORANGE),
     ("Credit Calculation", "Service credits and penalties computed automatically per contract terms", AMBER)],
    "Key Capabilities",
    [("9 SLA Metric Types",
      "Uptime, Response Time, Resolution Time, Throughput, Error Rate, Availability, Latency, MTTR, "
      "and Custom metrics \u2014 each with targets, thresholds, and measurement windows extracted by AI.", BLUE),
     ("Service Credit Recovery",
      "Automatic calculation of credits owed based on breach severity and contract terms. Track "
      "recovery status from detection through claim to resolution. Customers report 3\u20135x more credits.", EMERALD),
     ("Vendor Scorecards",
      "A\u2013F grade vendor ratings combining SLA compliance rate, breach frequency, financial impact, "
      "and trend direction. Compare vendor performance across your portfolio.", ORANGE)])

deep_dive(10, PINK,
    "Relationship Governance: only in Evaluetor",
    "Most CLM platforms stop at the contract. Evaluetor goes further \u2014 automatically connecting "
    "contract data to the business relationships they support. No other CLM vendor offers this capability.",
    "The Automated Governance Bridge", "Contract upload triggers a fully automated pipeline:",
    [("Contract Upload", "AI extracts parties, terms, SLAs, and obligations", ORANGE),
     ("Organizations", "Counterparties auto-linked to your org directory", BLUE),
     ("Relationships", "Business relationships created with health scoring", EMERALD),
     ("KPIs & Scoring", "Performance metrics tracked with dual-perspective scoring", AMBER),
     ("Improvement Points", "Gaps automatically surface with action recommendations", PINK)],
    "What Makes This Unique",
    [("Perception Gap Analysis",
      "Dual-perspective scoring reveals misalignment between how you see vendor performance and "
      "how they see it \u2014 with severity classification that drives targeted improvement.", PINK),
     ("Composite Health Scoring",
      "Relationship health calculated from Risk (30%), SLA Compliance (40%), and Obligation "
      "Fulfillment (30%) \u2014 a single score that tells you where each relationship truly stands.", EMERALD),
     ("Automated Improvement Tracking",
      "When perception gaps exceed thresholds, improvement points are auto-generated with ownership, "
      "priority, and action items. From insight to action without manual intervention.", ORANGE)])


# ================================================================
# SLIDE 11: MEASURABLE IMPACT
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
   "Measurable impact from day one", sz=28, clr=HEADING, bold=True)

# VALUE banner
rect(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.1), CARD)
tx(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
   "Immediate, quantifiable results \u2014 not months of implementation before you see",
   sz=14, clr=BODY, align=PP_ALIGN.CENTER)
tx(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.55),
   "VALUE.", sz=32, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)

# 4 metrics
metrics = [
    ("90%", "Less time\nreviewing contracts", PURPLE),
    ("3\u20135x", "More SLA credits\nrecovered", ORANGE),
    ("95%", "Fewer missed\ndeadlines", EMERALD),
    ("$2M+", "Credits recovered\nacross customers", BLUE),
]
for i, (val, label, accent) in enumerate(metrics):
    x = Inches(0.5) + Inches(i * 3.15)
    rect(s, x, Inches(2.2), Inches(2.95), Inches(1.5), CARD)
    tx(s, x + Inches(0.1), Inches(2.3), Inches(2.75), Inches(0.7),
       val, sz=38, clr=accent, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), Inches(3.0), Inches(2.75), Inches(0.55),
       label, sz=12, clr=BODY, align=PP_ALIGN.CENTER)

# 3 team cards
teams = [
    ("For Legal", "Days of contract review compressed to minutes. Automated risk scoring "
     "across 10 categories. Portfolio-wide semantic search.", PURPLE),
    ("For Procurement", "Real-time SLA monitoring. Detect breaches the moment they happen. "
     "Recover every service credit owed.", ORANGE),
    ("For Finance", "Track every obligation and deadline. Detect auto-renewal traps before "
     "they trigger. Understand true financial exposure.", EMERALD),
]
for i, (title, desc, accent) in enumerate(teams):
    x = Inches(0.5) + Inches(i * 4.2)
    rect(s, x, Inches(3.95), Inches(3.95), Inches(2.7), CARD)
    dot(s, x + Inches(0.2), Inches(4.15), accent, Inches(0.12))
    tx(s, x + Inches(0.45), Inches(4.07), Inches(3.3), Inches(0.3),
       title, sz=14, clr=HEADING, bold=True)
    tx(s, x + Inches(0.2), Inches(4.45), Inches(3.55), Inches(2.0),
       desc, sz=12, clr=BODY)

footer(s); sn(s, 11)


# ================================================================
# SLIDE 12: VALUE DELIVERED
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.15), Inches(3.2), Inches(0.4),
   "Evaluetor\u00a9 AI", sz=20, clr=BLUE, bold=True)
tx(s, Inches(3.5), Inches(0.15), Inches(6), Inches(0.4),
   "Value delivered", sz=22, clr=HEADING, bold=True)

values = [
    ("Increase return on contract\nbase by 3\u20136%",
     "Developed by seasoned industry experts"),
    ("Increase renewal rate and\ncontract extensions",
     "Trusted, transparent relationship \u2014\navoid transition costs"),
    ("Improve margins\n(decrease costs)",
     "Pro-actively action commitments\nand avoid omissions"),
    ("Increase revenue",
     "Growth of add-on business within\nthe existing customer portfolio"),
]
for i, (left, right) in enumerate(values):
    y = Inches(0.8) + Inches(i * 1.4)
    rect(s, Inches(0.5), y, Inches(5.0), Inches(1.2), CARD)
    tx(s, Inches(0.8), y + Inches(0.15), Inches(4.4), Inches(0.9),
       left, sz=15, clr=HEADING, bold=True, align=PP_ALIGN.CENTER)
    rect(s, Inches(5.8), y, Inches(7.0), Inches(1.2), CARD)
    tx(s, Inches(6.1), y + Inches(0.15), Inches(6.4), Inches(0.9),
       right, sz=15, clr=BODY, align=PP_ALIGN.CENTER)

# CTA button
cta = rect(s, Inches(4.0), Inches(6.45), Inches(5.3), Inches(0.5), CARD)
cta.fill.solid(); cta.fill.fore_color.rgb = ORANGE; cta.line.fill.background()
tx(s, Inches(4.0), Inches(6.47), Inches(5.3), Inches(0.45),
   "Let\u2019s schedule a live demo  \u2192", sz=17, clr=WHITE, bold=True,
   align=PP_ALIGN.CENTER)

footer(s); sn(s, 12)


# ================================================================
# SLIDE 13: COMPETITIVE COMPARISON — AI-NATIVE VS LEGACY CLM
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
   "AI-native. Not AI-bolted-on.", sz=28, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
   "Legacy CLM tools added AI as an afterthought. We built intelligence into the foundation.",
   sz=13, clr=BODY)

# Table: 5 columns × 9 rows (1 header + 8 capabilities)
COLS = [
    ("Capability",            Inches(3.2),  HEADING),
    ("Evaluetor",             Inches(2.3),  BLUE),
    ("DocuSign CLM",          Inches(2.3),  MUTED),
    ("Icertis",               Inches(2.3),  MUTED),
    ("Ironclad",              Inches(2.3),  MUTED),
]
ROWS = [
    ("Time to Value",            "Hours",          "6–12 months", "12–18 months", "3–6 months"),
    ("AI Extraction",            "9 native agents","Limited",     "Add-on",       "Limited"),
    ("SLA Monitoring",           "Built-in",       "No",          "Add-on",       "No"),
    ("Obligation Tracking",      "AI-extracted",   "Manual",      "Manual",       "Limited"),
    ("Relationship Governance",  "Built-in",       "No",          "No",           "No"),
    ("Knowledge Graph",          "Built-in",       "No",          "No",           "No"),
    ("Unlimited Users",          "✓ Yes",          "✗",           "✗",            "✗"),
    ("LLM Observability",        "Langfuse",       "No",          "No",           "No"),
]

# Header row
x0 = Inches(0.5); y = Inches(1.45)
row_h = Inches(0.45)
header_h = Inches(0.5)
hdr = rect(s, x0, y, Inches(12.4), header_h, HEADING)
hdr.fill.solid(); hdr.fill.fore_color.rgb = HEADING
x = x0
for (label, w, _clr) in COLS:
    tx(s, x + Inches(0.15), y + Inches(0.1), w - Inches(0.2), Inches(0.35),
       label, sz=12, clr=WHITE, bold=True,
       align=PP_ALIGN.LEFT if label == "Capability" else PP_ALIGN.CENTER)
    x += w

# Data rows
for ri, row in enumerate(ROWS):
    ry = y + header_h + Inches(ri * 0.5)
    fill = CARD if ri % 2 == 0 else CARD_ALT
    rect(s, x0, ry, Inches(12.4), row_h, fill)
    x = x0
    for ci, val in enumerate(row):
        is_evaluetor = (ci == 1)
        clr = BLUE if is_evaluetor else (HEADING if ci == 0 else BODY)
        bold = is_evaluetor or ci == 0
        col_w = COLS[ci][1]
        tx(s, x + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.2), Inches(0.32),
           val, sz=11, clr=clr, bold=bold,
           align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
        x += col_w

# Bottom tagline
rect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.7), RGBColor(0xEF, 0xF6, 0xFF))
tx(s, Inches(0.8), Inches(6.18), Inches(11.7), Inches(0.45),
   "Time-to-value in hours, not months — because the intelligence is built in, not bolted on.",
   sz=14, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)

footer(s); sn(s, 13)


# ================================================================
# SLIDE 14: CONTACT
# ================================================================
s = prs.slides.add_slide(BL); set_bg(s); bar(s)
tx(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
   "Let’s talk", sz=28, clr=HEADING, bold=True)
tx(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
   "Ready to see what’s really inside your contract portfolio?",
   sz=14, clr=BODY)

# Main contact card
rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0), CARD)

# Left side — person
tx(s, Inches(1.0), Inches(1.9), Inches(7), Inches(0.5),
   "Steven Visser", sz=32, clr=HEADING, bold=True)
tx(s, Inches(1.0), Inches(2.6), Inches(7), Inches(0.35),
   "Managing Partner — Evaluetor© AI", sz=16, clr=BLUE, bold=True)

# Address
tx(s, Inches(1.0), Inches(3.4), Inches(7), Inches(0.3),
   "Burgemeester Le Fevre de Montignylaan 264", sz=13, clr=BODY)
tx(s, Inches(1.0), Inches(3.75), Inches(7), Inches(0.3),
   "3055 NK Rotterdam, The Netherlands", sz=13, clr=BODY)

# Contact details with accents
details = [
    ("Mobile",  "+31 (0) 6 31003495",      ORANGE),
    ("Email",   "Steven.Visser@Evaluetor.com", BLUE),
    ("Web",     "www.Evaluetor.com",       EMERALD),
]
for i, (label, val, accent) in enumerate(details):
    y = Inches(4.45) + Inches(i * 0.55)
    dot(s, Inches(1.0), y + Inches(0.12), accent, Inches(0.14))
    tx(s, Inches(1.3), y, Inches(1.2), Inches(0.4),
       label, sz=12, clr=MUTED, bold=True)
    tx(s, Inches(2.5), y, Inches(6.0), Inches(0.4),
       val, sz=14, clr=HEADING, bold=True)

# Right side — CTA panel
rect(s, Inches(8.5), Inches(2.0), Inches(4.0), Inches(4.0), CARD_ALT)
tx(s, Inches(8.7), Inches(2.3), Inches(3.6), Inches(0.4),
   "Evaluetor© AI", sz=18, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(8.7), Inches(2.75), Inches(3.6), Inches(0.35),
   "AI-Powered Contract Intelligence", sz=11, clr=MUTED,
   bold=True, align=PP_ALIGN.CENTER)

# Mini stats
mini = [("9", "AI agents"), ("30+", "Clause types"), ("10", "Risk categories")]
for i, (val, label) in enumerate(mini):
    y = Inches(3.4) + Inches(i * 0.75)
    tx(s, Inches(8.7), y, Inches(3.6), Inches(0.4),
       val, sz=22, clr=BLUE, bold=True, align=PP_ALIGN.CENTER)
    tx(s, Inches(8.7), y + Inches(0.4), Inches(3.6), Inches(0.3),
       label, sz=10, clr=BODY, align=PP_ALIGN.CENTER)

footer(s); sn(s, 14)


# ─── Save ───
out = "/Users/jjayaraj/workspaces/studios/clm/sales/Evaluetor-Commercial-Presentation.pptx"
prs.save(out); print(f"Saved: {out}")
