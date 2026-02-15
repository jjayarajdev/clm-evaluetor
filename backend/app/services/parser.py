"""Document parsing service for extracting text from PDFs and DOCX files."""

import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import BinaryIO

import fitz  # PyMuPDF
import xlrd
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

logger = logging.getLogger(__name__)


# Minimum characters per page to consider it has valid text (not scanned)
MIN_TEXT_THRESHOLD = 50


@dataclass
class PageContent:
    """Content from a single page."""

    page_number: int
    text: str
    char_count: int = 0

    def __post_init__(self) -> None:
        self.char_count = len(self.text)


@dataclass
class DocumentMetadata:
    """Metadata extracted from the document."""

    title: str | None = None
    author: str | None = None
    subject: str | None = None
    creator: str | None = None
    creation_date: str | None = None
    modification_date: str | None = None
    page_count: int = 0
    word_count: int = 0


@dataclass
class ParsedDocument:
    """Result of parsing a document."""

    filename: str
    file_type: str  # "pdf" or "docx"
    pages: list[PageContent] = field(default_factory=list)
    full_text: str = ""
    metadata: DocumentMetadata = field(default_factory=DocumentMetadata)
    ocr_used: bool = False
    success: bool = True
    error: str | None = None

    @property
    def page_count(self) -> int:
        """Get total page count."""
        return len(self.pages)

    @property
    def total_chars(self) -> int:
        """Get total character count."""
        return sum(p.char_count for p in self.pages)


class ParserError(Exception):
    """Exception raised for parsing errors."""

    pass


class DocumentParser:
    """Service for parsing PDF and DOCX documents."""

    def __init__(self, enable_ocr: bool = True) -> None:
        """Initialize parser.

        Args:
            enable_ocr: Whether to use OCR for scanned PDFs.
        """
        self.enable_ocr = enable_ocr
        self._tesseract_available: bool | None = None

    def check_tesseract(self) -> bool:
        """Check if Tesseract is available."""
        if self._tesseract_available is not None:
            return self._tesseract_available

        try:
            import pytesseract

            pytesseract.get_tesseract_version()
            self._tesseract_available = True
        except Exception:
            logger.warning("Tesseract not available - OCR will be disabled")
            self._tesseract_available = False

        return self._tesseract_available

    def parse_file(self, file_path: str | Path) -> ParsedDocument:
        """Parse a document from a file path.

        Args:
            file_path: Path to the document.

        Returns:
            ParsedDocument with extracted content.
        """
        path = Path(file_path)
        if not path.exists():
            return ParsedDocument(
                filename=path.name,
                file_type="unknown",
                success=False,
                error=f"File not found: {file_path}",
            )

        ext = path.suffix.lower()

        with open(path, "rb") as f:
            if ext == ".pdf":
                return self.parse_pdf(f, path.name)
            elif ext == ".docx":
                return self.parse_docx(f, path.name)
            elif ext == ".doc":
                # Older Word format - try to extract with limited support
                return self._parse_doc_legacy(path, path.name)
            elif ext == ".xlsx":
                return self.parse_excel(f, path.name)
            elif ext == ".xls":
                # Legacy Excel format - use xlrd
                return self._parse_xls_legacy(path, path.name)
            elif ext == ".pptx":
                return self.parse_powerpoint(f, path.name)
            elif ext == ".ppt":
                # Legacy PowerPoint format - binary extraction
                return self._parse_ppt_legacy(path, path.name)
            else:
                return ParsedDocument(
                    filename=path.name,
                    file_type=ext,
                    success=False,
                    error=f"Unsupported file type: {ext}",
                )

    def parse_pdf(self, file: BinaryIO, filename: str = "document.pdf") -> ParsedDocument:
        """Parse a PDF document.

        Args:
            file: File-like object containing PDF data.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            # Read the entire content
            content = file.read()
            doc = fitz.open(stream=content, filetype="pdf")

            pages: list[PageContent] = []
            all_text_parts: list[str] = []
            needs_ocr = False

            # Extract text from each page
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text("text")

                # Clean up text
                text = text.strip()

                # Check if page needs OCR (minimal text extracted)
                if len(text) < MIN_TEXT_THRESHOLD:
                    needs_ocr = True

                pages.append(
                    PageContent(
                        page_number=page_num + 1,
                        text=text,
                    )
                )
                if text:
                    all_text_parts.append(f"[Page {page_num + 1}]\n{text}")

            # Try OCR if needed and available
            ocr_used = False
            if needs_ocr and self.enable_ocr and self.check_tesseract():
                logger.info(f"Attempting OCR for {filename}")
                ocr_pages = self._ocr_pdf(doc)
                if ocr_pages:
                    # Merge OCR results with existing pages
                    pages = self._merge_ocr_results(pages, ocr_pages)
                    all_text_parts = [
                        f"[Page {p.page_number}]\n{p.text}" for p in pages if p.text
                    ]
                    ocr_used = True

            # Extract metadata
            metadata = self._extract_pdf_metadata(doc)
            doc.close()

            # Combine all text
            full_text = "\n\n".join(all_text_parts)
            metadata.word_count = len(full_text.split())

            return ParsedDocument(
                filename=filename,
                file_type="pdf",
                pages=pages,
                full_text=full_text,
                metadata=metadata,
                ocr_used=ocr_used,
                success=True,
            )

        except Exception as e:
            logger.exception(f"Error parsing PDF {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="pdf",
                success=False,
                error=str(e),
            )

    def parse_docx(self, file: BinaryIO, filename: str = "document.docx") -> ParsedDocument:
        """Parse a DOCX document.

        Args:
            file: File-like object containing DOCX data.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            doc = Document(file)

            # Extract paragraphs with structure preservation
            paragraphs: list[str] = []
            current_page = 1
            page_contents: dict[int, list[str]] = {1: []}

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect heading styles
                style_name = para.style.name if para.style else ""
                if style_name.startswith("Heading"):
                    # Add heading marker
                    level = style_name.replace("Heading ", "").replace("Heading", "1")
                    try:
                        level_num = int(level)
                    except ValueError:
                        level_num = 1
                    prefix = "#" * level_num + " "
                    text = f"\n{prefix}{text}\n"

                paragraphs.append(text)

                # DOCX doesn't have page numbers in the same way as PDF
                # We'll estimate pages based on character count (approx 3000 chars/page)
                total_chars = sum(len(p) for p in page_contents[current_page])
                if total_chars > 3000:
                    current_page += 1
                    page_contents[current_page] = []
                page_contents[current_page].append(text)

            # Build pages
            pages = []
            for page_num in sorted(page_contents.keys()):
                page_text = "\n".join(page_contents[page_num])
                if page_text.strip():
                    pages.append(
                        PageContent(
                            page_number=page_num,
                            text=page_text,
                        )
                    )

            # Extract tables
            table_texts = []
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    table_texts.append(table_text)

            # Combine all content
            full_text = "\n\n".join(paragraphs)
            if table_texts:
                full_text += "\n\n[Tables]\n" + "\n\n".join(table_texts)

            # Extract metadata from document properties
            metadata = self._extract_docx_metadata(doc)
            metadata.page_count = len(pages)
            metadata.word_count = len(full_text.split())

            return ParsedDocument(
                filename=filename,
                file_type="docx",
                pages=pages,
                full_text=full_text,
                metadata=metadata,
                ocr_used=False,
                success=True,
            )

        except PackageNotFoundError:
            return ParsedDocument(
                filename=filename,
                file_type="docx",
                success=False,
                error="Invalid or corrupted DOCX file",
            )
        except Exception as e:
            logger.exception(f"Error parsing DOCX {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="docx",
                success=False,
                error=str(e),
            )

    def parse_excel(self, file: BinaryIO, filename: str = "document.xlsx") -> ParsedDocument:
        """Parse an Excel document.

        Args:
            file: File-like object containing Excel data.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            wb = load_workbook(file, read_only=True, data_only=True)

            pages: list[PageContent] = []
            all_text_parts: list[str] = []

            # Each sheet becomes a "page"
            for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
                ws = wb[sheet_name]
                sheet_text_parts = [f"## Sheet: {sheet_name}\n"]

                # Extract data from cells
                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell).strip())
                            else:
                                row_text.append("")
                        rows_data.append(row_text)

                # Format as table
                if rows_data:
                    # Use first row as headers if it looks like headers
                    if len(rows_data) > 1:
                        headers = rows_data[0]
                        sheet_text_parts.append(" | ".join(headers))
                        sheet_text_parts.append("-" * 50)
                        for row in rows_data[1:]:
                            # Pad row to match header length
                            while len(row) < len(headers):
                                row.append("")
                            sheet_text_parts.append(" | ".join(row[:len(headers)]))
                    else:
                        for row in rows_data:
                            sheet_text_parts.append(" | ".join(row))

                sheet_text = "\n".join(sheet_text_parts)
                pages.append(
                    PageContent(
                        page_number=sheet_num,
                        text=sheet_text,
                    )
                )
                all_text_parts.append(f"[Sheet {sheet_num}: {sheet_name}]\n{sheet_text}")

            wb.close()

            # Combine all text
            full_text = "\n\n".join(all_text_parts)

            metadata = DocumentMetadata(
                title=filename,
                page_count=len(pages),
                word_count=len(full_text.split()),
            )

            return ParsedDocument(
                filename=filename,
                file_type="xlsx",
                pages=pages,
                full_text=full_text,
                metadata=metadata,
                ocr_used=False,
                success=True,
            )

        except InvalidFileException:
            return ParsedDocument(
                filename=filename,
                file_type="xlsx",
                success=False,
                error="Invalid or corrupted Excel file",
            )
        except Exception as e:
            logger.exception(f"Error parsing Excel {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="xlsx",
                success=False,
                error=str(e),
            )

    def parse_powerpoint(self, file: BinaryIO, filename: str = "document.pptx") -> ParsedDocument:
        """Parse a PowerPoint document.

        Args:
            file: File-like object containing PowerPoint data.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            prs = Presentation(file)

            pages: list[PageContent] = []
            all_text_parts: list[str] = []

            # Each slide becomes a "page"
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = [f"## Slide {slide_num}\n"]

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)

                    # Extract text from tables
                    if shape.has_table:
                        table_text = self._extract_pptx_table(shape.table)
                        if table_text:
                            slide_text_parts.append(table_text)

                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    pages.append(
                        PageContent(
                            page_number=slide_num,
                            text=slide_text,
                        )
                    )
                    all_text_parts.append(f"[Slide {slide_num}]\n{slide_text}")

            # Combine all text
            full_text = "\n\n".join(all_text_parts)

            metadata = DocumentMetadata(
                title=filename,
                page_count=len(pages),
                word_count=len(full_text.split()),
            )

            return ParsedDocument(
                filename=filename,
                file_type="pptx",
                pages=pages,
                full_text=full_text,
                metadata=metadata,
                ocr_used=False,
                success=True,
            )

        except PptxPackageNotFoundError:
            return ParsedDocument(
                filename=filename,
                file_type="pptx",
                success=False,
                error="Invalid or corrupted PowerPoint file",
            )
        except Exception as e:
            logger.exception(f"Error parsing PowerPoint {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="pptx",
                success=False,
                error=str(e),
            )

    def _extract_pptx_table(self, table) -> str:
        """Extract text from a PowerPoint table."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _parse_doc_legacy(self, file_path: Path, filename: str) -> ParsedDocument:
        """Parse older .doc Word format using textract or fallback.

        Args:
            file_path: Path to the .doc file.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            # Try using antiword if available
            import subprocess
            result = subprocess.run(
                ["antiword", str(file_path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout:
                full_text = result.stdout.strip()

                # Create a single page with all content
                pages = [PageContent(page_number=1, text=full_text)]

                metadata = DocumentMetadata(
                    title=filename,
                    page_count=1,
                    word_count=len(full_text.split()),
                )

                return ParsedDocument(
                    filename=filename,
                    file_type="doc",
                    pages=pages,
                    full_text=full_text,
                    metadata=metadata,
                    ocr_used=False,
                    success=True,
                )
        except FileNotFoundError:
            logger.warning("antiword not installed - .doc files may not be fully supported")
        except subprocess.TimeoutExpired:
            logger.warning(f"Timeout parsing .doc file: {filename}")
        except Exception as e:
            logger.warning(f"Error with antiword for {filename}: {e}")

        # Fallback: try to read as binary and extract text
        try:
            with open(file_path, "rb") as f:
                content = f.read()
                # Simple extraction of readable text from binary
                text_parts = []
                current_text = []
                for byte in content:
                    if 32 <= byte <= 126:  # Printable ASCII
                        current_text.append(chr(byte))
                    else:
                        if len(current_text) > 10:  # Only keep meaningful strings
                            text_parts.append("".join(current_text))
                        current_text = []
                if current_text and len(current_text) > 10:
                    text_parts.append("".join(current_text))

                full_text = " ".join(text_parts)

                if len(full_text) > 100:  # Only if we got meaningful content
                    pages = [PageContent(page_number=1, text=full_text)]
                    metadata = DocumentMetadata(
                        title=filename,
                        page_count=1,
                        word_count=len(full_text.split()),
                    )

                    return ParsedDocument(
                        filename=filename,
                        file_type="doc",
                        pages=pages,
                        full_text=full_text,
                        metadata=metadata,
                        ocr_used=False,
                        success=True,
                    )
        except Exception as e:
            logger.warning(f"Fallback extraction failed for {filename}: {e}")

        return ParsedDocument(
            filename=filename,
            file_type="doc",
            success=False,
            error="Unable to parse .doc file. Please convert to .docx format for better results.",
        )

    def _parse_xls_legacy(self, file_path: Path, filename: str) -> ParsedDocument:
        """Parse older .xls Excel format using xlrd.

        Args:
            file_path: Path to the .xls file.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            wb = xlrd.open_workbook(str(file_path))

            pages: list[PageContent] = []
            all_text_parts: list[str] = []

            # Each sheet becomes a "page"
            for sheet_num, sheet_name in enumerate(wb.sheet_names(), 1):
                ws = wb.sheet_by_name(sheet_name)
                sheet_text_parts = [f"## Sheet: {sheet_name}\n"]

                rows_data = []
                for row_idx in range(ws.nrows):
                    row = ws.row_values(row_idx)
                    # Filter out completely empty rows
                    if any(cell for cell in row):
                        row_text = []
                        for cell in row:
                            if cell:
                                row_text.append(str(cell).strip())
                            else:
                                row_text.append("")
                        rows_data.append(row_text)

                # Format as table
                if rows_data:
                    if len(rows_data) > 1:
                        headers = rows_data[0]
                        sheet_text_parts.append(" | ".join(str(h) for h in headers))
                        sheet_text_parts.append("-" * 50)
                        for row in rows_data[1:]:
                            while len(row) < len(headers):
                                row.append("")
                            sheet_text_parts.append(" | ".join(str(c) for c in row[:len(headers)]))
                    else:
                        for row in rows_data:
                            sheet_text_parts.append(" | ".join(str(c) for c in row))

                sheet_text = "\n".join(sheet_text_parts)
                pages.append(
                    PageContent(
                        page_number=sheet_num,
                        text=sheet_text,
                    )
                )
                all_text_parts.append(f"[Sheet {sheet_num}: {sheet_name}]\n{sheet_text}")

            # Combine all text
            full_text = "\n\n".join(all_text_parts)

            metadata = DocumentMetadata(
                title=filename,
                page_count=len(pages),
                word_count=len(full_text.split()),
            )

            return ParsedDocument(
                filename=filename,
                file_type="xls",
                pages=pages,
                full_text=full_text,
                metadata=metadata,
                ocr_used=False,
                success=True,
            )

        except xlrd.biffh.XLRDError as e:
            logger.exception(f"Error parsing legacy Excel {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="xls",
                success=False,
                error=f"Invalid or corrupted .xls file: {e}",
            )
        except Exception as e:
            logger.exception(f"Error parsing legacy Excel {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="xls",
                success=False,
                error=str(e),
            )

    def _parse_ppt_legacy(self, file_path: Path, filename: str) -> ParsedDocument:
        """Parse older .ppt PowerPoint format using binary extraction.

        Args:
            file_path: Path to the .ppt file.
            filename: Original filename for reference.

        Returns:
            ParsedDocument with extracted content.
        """
        try:
            with open(file_path, "rb") as f:
                content = f.read()

            # Extract readable text from the binary OLE format
            # PPT files contain text in various places, we'll extract printable strings
            text_parts = []
            current_text = []

            for byte in content:
                if 32 <= byte <= 126:  # Printable ASCII
                    current_text.append(chr(byte))
                elif byte in (9, 10, 13):  # Tab, newline, carriage return
                    current_text.append(' ')
                else:
                    if len(current_text) > 5:  # Keep meaningful strings
                        text = "".join(current_text).strip()
                        # Filter out common binary garbage patterns
                        if text and not self._is_binary_garbage(text):
                            text_parts.append(text)
                    current_text = []

            if current_text and len(current_text) > 5:
                text = "".join(current_text).strip()
                if text and not self._is_binary_garbage(text):
                    text_parts.append(text)

            # Deduplicate while preserving order and filter very short fragments
            seen = set()
            unique_parts = []
            for part in text_parts:
                if len(part) > 10 and part not in seen:
                    seen.add(part)
                    unique_parts.append(part)

            full_text = "\n".join(unique_parts)

            if len(full_text) > 50:  # Only if we got meaningful content
                pages = [PageContent(page_number=1, text=full_text)]
                metadata = DocumentMetadata(
                    title=filename,
                    page_count=1,
                    word_count=len(full_text.split()),
                )

                return ParsedDocument(
                    filename=filename,
                    file_type="ppt",
                    pages=pages,
                    full_text=full_text,
                    metadata=metadata,
                    ocr_used=False,
                    success=True,
                )

            return ParsedDocument(
                filename=filename,
                file_type="ppt",
                success=False,
                error="Could not extract text from .ppt file. Please convert to .pptx format for better results.",
            )

        except Exception as e:
            logger.exception(f"Error parsing legacy PowerPoint {filename}")
            return ParsedDocument(
                filename=filename,
                file_type="ppt",
                success=False,
                error=str(e),
            )

    def _is_binary_garbage(self, text: str) -> bool:
        """Check if extracted text looks like binary garbage."""
        if not text:
            return True

        # Count alphanumeric vs special characters
        alpha_count = sum(1 for c in text if c.isalnum() or c.isspace())
        total = len(text)

        # If less than 50% alphanumeric, probably garbage
        if total > 0 and alpha_count / total < 0.5:
            return True

        # Check for repetitive patterns that indicate binary data
        if len(set(text)) < len(text) * 0.1:  # Very few unique chars
            return True

        return False

    def _extract_pdf_metadata(self, doc: fitz.Document) -> DocumentMetadata:
        """Extract metadata from PDF document."""
        meta = doc.metadata or {}

        return DocumentMetadata(
            title=meta.get("title") or None,
            author=meta.get("author") or None,
            subject=meta.get("subject") or None,
            creator=meta.get("creator") or None,
            creation_date=meta.get("creationDate") or None,
            modification_date=meta.get("modDate") or None,
            page_count=len(doc),
        )

    def _extract_docx_metadata(self, doc: Document) -> DocumentMetadata:
        """Extract metadata from DOCX document."""
        props = doc.core_properties

        creation_date = None
        modification_date = None

        if props.created:
            creation_date = props.created.isoformat()
        if props.modified:
            modification_date = props.modified.isoformat()

        return DocumentMetadata(
            title=props.title or None,
            author=props.author or None,
            subject=props.subject or None,
            creator=props.author or None,  # DOCX uses author for creator
            creation_date=creation_date,
            modification_date=modification_date,
        )

    def _extract_table_text(self, table) -> str:
        """Extract text from a DOCX table."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _ocr_pdf(self, doc: fitz.Document) -> list[PageContent]:
        """Perform OCR on PDF pages using Tesseract.

        Args:
            doc: PyMuPDF document object.

        Returns:
            List of PageContent with OCR text.
        """
        import pytesseract
        from PIL import Image

        pages = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Render page to image at 300 DPI for better OCR
            mat = fitz.Matrix(300 / 72, 300 / 72)  # 300 DPI
            pix = page.get_pixmap(matrix=mat)

            # Convert to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Run OCR
            try:
                text = pytesseract.image_to_string(img, lang="eng")
                text = text.strip()
            except Exception as e:
                logger.warning(f"OCR failed for page {page_num + 1}: {e}")
                text = ""

            pages.append(
                PageContent(
                    page_number=page_num + 1,
                    text=text,
                )
            )

        return pages

    def _merge_ocr_results(
        self, original: list[PageContent], ocr: list[PageContent]
    ) -> list[PageContent]:
        """Merge OCR results with original text extraction.

        Uses OCR text only if it has more content than original.

        Args:
            original: Original text extraction results.
            ocr: OCR extraction results.

        Returns:
            Merged list of PageContent.
        """
        merged = []

        for orig, ocr_page in zip(original, ocr):
            # Use whichever has more text
            if ocr_page.char_count > orig.char_count:
                merged.append(ocr_page)
            else:
                merged.append(orig)

        return merged


# Singleton instance
_parser: DocumentParser | None = None


def get_parser() -> DocumentParser:
    """Get the document parser singleton."""
    global _parser
    if _parser is None:
        _parser = DocumentParser()
    return _parser
